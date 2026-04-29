#!/usr/bin/env python3
"""
Edit Sentinel_AI_Presentation - Repaired.pptx
Adds: Purpose, Hardware Photos, Timeline, Team Contributions,
      Hardware Constraints, Software Constraints, Future Work
Reorders all slides into logical presentation flow.
Output: /Users/karthi/Desktop/Sentinal_AI/Sentinel_AI_Final.pptx
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# ── Exact design tokens matched from repaired PPTX ────────────────────────────
C_NAV  = RGBColor(0x1E, 0x2A, 0x3A)   # header/footer bars
C_WHITE= RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE = RGBColor(0x41, 0x6E, 0x8F)   # labels, accents
C_DARK = RGBColor(0x2C, 0x3E, 0x50)   # body text
C_GREEN= RGBColor(0x1A, 0x9E, 0x8A)
C_GRAY = RGBColor(0x7F, 0x8C, 0x8D)   # footer text
C_CARD = RGBColor(0xF4, 0xF6, 0xF8)   # card bg

# ── Exact measurements from template (EMU) ────────────────────────────────────
SW   = Emu(12188952)
HDR_H= Emu(1005840)
FTR_T= Emu(6492240)
FTR_H= Emu(365760)
# Title text box (inside header)
TTL_L= Emu(274320);  TTL_T= Emu(137160);  TTL_W= Emu(11430000); TTL_H= Emu(492443)
# Footer text box
FTX_L= Emu(182880);  FTX_T= Emu(6510528); FTX_W= Emu(11887200); FTX_H= Emu(320040)
# Left-column card pattern (from slide 2):
#   3 stacked cards on left, stats row on right
LC_L = Emu(274320);   LC_W = Emu(5029200);  LC_H = Emu(1371600)
LC_TOPS = [Emu(1143000), Emu(2697480), Emu(4251960)]
LBL_L= Emu(411480);  LBL_H= Emu(261610)   # label box inside card
LBL_TOPS=[Emu(1216152), Emu(2770632), Emu(4325112)]
BDY_H= Emu(914400)                          # body text height
BDY_TOPS=[Emu(1554480), Emu(3108960), Emu(4663440)]
# Stat boxes (right side of slide 2):
ST_TOPS= Emu(4023360); ST_H= Emu(1234440)
ST_LS  = [Emu(5760720), Emu(7333488), Emu(8906256), Emu(10479024)]
ST_W   = Emu(1417320)
# 5-card row pattern (from slide 7 — How Monitors Work):
C5_T = Emu(3931920); C5_H = Emu(1554480); C5_W = Emu(2240280)
C5_LS= [Emu(228600), Emu(2587752), Emu(4946904), Emu(7306056), Emu(9665208)]
# 3-col table pattern (from slide 5 — Challenges):
TBL_ROW_H = Emu(256848)
TBL_HDR_H = Emu(265608)
TBL_HDR_T = Emu(1124496)
TBL_C1_L=Emu(320040); TBL_C1_W=Emu(3108560)
TBL_C2_L=Emu(3520040); TBL_C2_W=Emu(3508560)
TBL_C3_L=Emu(7120040); TBL_C3_W=Emu(4748872)


# ── Core helpers ──────────────────────────────────────────────────────────────
def _rect(sl, l, t, w, h, rgb):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = rgb
    s.line.fill.background()
    return s

def _tb(sl, l, t, w, h, text, sz_emu=114300, bold=False, color=None,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Emu(sz_emu); r.font.bold = bold
    if color: r.font.color.rgb = color
    return tb

def _header(sl, title):
    _rect(sl, Emu(0), Emu(0), SW, HDR_H, C_NAV)
    _tb(sl, TTL_L, TTL_T, TTL_W, TTL_H, title,
        sz_emu=330200, bold=True, color=C_WHITE)

def _footer(sl, ref=""):
    _rect(sl, Emu(0), FTR_T, SW, FTR_H, C_NAV)
    if ref:
        _tb(sl, FTX_L, FTX_T, FTX_W, FTX_H, ref, sz_emu=101600, color=C_GRAY)

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ── Left-card + stats layout (matches slide 2 exactly) ───────────────────────
def _left_card(sl, idx, label, body):
    """Add one of the 3 stacked left cards (idx=0,1,2)."""
    _rect(sl, LC_L, LC_TOPS[idx], LC_W, LC_H, C_CARD)
    _tb(sl, LBL_L, LBL_TOPS[idx], Emu(4754880), LBL_H,
        label, sz_emu=139700, bold=True, color=C_BLUE)
    _tb(sl, LBL_L, BDY_TOPS[idx], Emu(4754880), BDY_H,
        body, sz_emu=120650, color=C_DARK)

def _stat_box(sl, idx, number, label):
    """One of 4 stat boxes on the right (idx=0..3)."""
    _rect(sl, ST_LS[idx], ST_TOPS, ST_W, ST_H, C_CARD)
    _tb(sl, ST_LS[idx]+Emu(91440), ST_TOPS+Emu(91440), Emu(1234440), Emu(502920),
        number, sz_emu=279400, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    _tb(sl, ST_LS[idx]+Emu(91440), ST_TOPS+Emu(640080), Emu(1234440), Emu(640080),
        label, sz_emu=114300, color=C_DARK, align=PP_ALIGN.CENTER)

# ── 3-col table helpers (matches slide 5 exactly) ────────────────────────────
def _tbl_header(sl, c1, c2, c3):
    for l, w, txt in [(TBL_C1_L,TBL_C1_W,c1),(TBL_C2_L,TBL_C2_W,c2),(TBL_C3_L,TBL_C3_W,c3)]:
        _rect(sl, l, TBL_HDR_T, w, TBL_HDR_H, C_BLUE)
        _tb(sl, l+Emu(45720), TBL_HDR_T+Emu(27216), w-Emu(91440), TBL_HDR_H-Emu(54432),
            txt, sz_emu=114300, bold=True, color=C_WHITE)

def _tbl_row(sl, row_num, c1, c2, c3):
    """row_num=0 is first data row."""
    t = TBL_HDR_T + TBL_HDR_H + TBL_ROW_H * row_num
    bg = C_CARD if row_num % 2 == 0 else C_WHITE
    for l, w, txt in [(TBL_C1_L,TBL_C1_W,c1),(TBL_C2_L,TBL_C2_W,c2),(TBL_C3_L,TBL_C3_W,c3)]:
        _rect(sl, l, t, w, TBL_ROW_H, bg)
        _tb(sl, l+Emu(45720), t+Emu(27216), w-Emu(91440), TBL_ROW_H-Emu(54432),
            txt, sz_emu=107950, color=C_DARK)

# ── 4-card row helper (similar to 5-card but 4 cards) ────────────────────────
def _four_cards(sl, top, height, cards):
    """4 equal-width cards across the content area. cards=[(label,body),...]"""
    w = Emu((12188952 - 228600*2 - 91440*3) // 4)
    gap = Emu(91440)
    l0 = Emu(228600)
    for i, (label, body) in enumerate(cards):
        l = l0 + (w + gap) * i
        _rect(sl, l, top, w, height, C_CARD)
        _tb(sl, l+Emu(91440), top+Emu(73152), w-Emu(182880), Emu(246221),
            label, sz_emu=127000, bold=True, color=C_BLUE)
        _tb(sl, l+Emu(91440), top+Emu(360000), w-Emu(182880), height-Emu(410000),
            body, sz_emu=107950, color=C_DARK)


# ══════════════════════════════════════════════════════════════════════════════
#  NEW SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def build_purpose(prs):
    """Matches the left-card + stats layout of slide 2 (Project Overview)."""
    sl = _blank(prs)
    _header(sl, "Purpose — Why Does This System Exist?")
    _footer(sl, "Temperature & humidity govern crop health, cold chains, industrial safety and patient care — unattended monitoring saves lives and yields")

    # 3 left cards (same geometry as slide 2)
    _left_card(sl, 0, "The Problem",
               "Crops, factories and hospitals rely on IoT sensors that fail "
               "silently. A cold snap destroys a harvest. A motor overheats in "
               "a factory. A vaccine fridge warms up overnight. Nobody notices "
               "until the damage is done.")

    _left_card(sl, 1, "Agriculture — Primary Use Case",
               "The AHT20 sensor tracks temperature (°C) and humidity (%) "
               "continuously. Sentinel AI detects any drift beyond the learned "
               "normal range within 10 seconds and triggers automated alerts "
               "and corrective actions — protecting crops without human supervision.")

    _left_card(sl, 2, "The System Solves This",
               "Multi-agent pipeline: Detect → Diagnose (LLM AI) → Alert "
               "(toast + log) → Recover (automated action) → Learn (baseline "
               "updated). Works across heterogeneous hardware — edge nodes, "
               "Pi bridges and cloud AI — all coordinated in real time.")

    # Right side: image area replaced with 4 application domain boxes
    app_top  = Emu(1143000)
    app_w    = Emu(1600000)
    app_h    = Emu(730000)
    app_gap  = Emu(80000)
    app_l    = Emu(5577840)

    apps = [
        ("Agriculture",     "Greenhouse climate,\ncrop frost alerts"),
        ("Industrial IoT",  "Motor temp, predictive\nmaintenance"),
        ("Healthcare",      "Cold-chain drugs,\nsterile room humidity"),
        ("Data Centres",    "Rack cooling,\nfire-risk humidity"),
        ("Smart Buildings", "HVAC optimisation,\nenergy saving"),
        ("Environment",     "Flood sensors,\nweather stations"),
    ]
    for i, (title, desc) in enumerate(apps):
        col = i % 2
        row = i // 2
        l = app_l + (app_w + app_gap) * col
        t = app_top + (app_h + app_gap) * row
        _rect(sl, l, t, app_w, app_h, C_CARD)
        _rect(sl, l, t, Emu(50000), app_h, C_BLUE)   # accent strip
        _tb(sl, l+Emu(100000), t+Emu(60000), app_w-Emu(130000), Emu(220000),
            title, sz_emu=127000, bold=True, color=C_BLUE)
        _tb(sl, l+Emu(100000), t+Emu(310000), app_w-Emu(130000), Emu(360000),
            desc, sz_emu=107950, color=C_DARK)

    # Stat row at bottom right (same positions as slide 2)
    _stat_box(sl, 0, "5s",    "collect\ninterval")
    _stat_box(sl, 1, "<10s",  "anomaly\ndetected")
    _stat_box(sl, 2, "6",     "AI\nagents")
    _stat_box(sl, 3, "15+",   "recovery\nactions")


def build_hardware_photos(prs):
    """3 real photos: AHT20, LoRa32, full setup + spec cards below."""
    sl = _blank(prs)
    _header(sl, "Hardware Used — Real Components")
    _footer(sl, "All three units connected and tested live  ·  LoRa32 → USB-C serial → Raspberry Pi 4B → Wi-Fi → Mac AI hub")

    IMG_AHT20  = "/tmp/hw_imgs/aht20.jpg"
    IMG_LORA32 = "/tmp/hw_imgs/lora32.jpg"
    IMG_SETUP  = "/tmp/hw_imgs/setup.jpg"

    gap = Emu(91440)
    pw  = (Emu(12188952) - Emu(228600)*2 - gap*2) // 3
    ph  = Emu(3300000)
    pt  = Emu(1143000)

    # Header labels ABOVE photos
    labels = [
        ("AHT20 — Temp & Humidity Sensor",       C_GREEN),
        ("Heltec WiFi LoRa32 V3 — Long-Range Node", C_BLUE),
        ("Full Setup: LoRa32 + OLED + Raspberry Pi 4B", RGBColor(0xE6,0x7E,0x22)),
    ]
    for i, (lbl, col) in enumerate(labels):
        lx = Emu(228600) + (pw+gap)*i
        _rect(sl, lx, pt-Emu(260000), pw, Emu(255000), col)
        _tb(sl, lx+Emu(45720), pt-Emu(230000), pw-Emu(91440), Emu(210000),
            lbl, sz_emu=114300, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Photos
    for i, img in enumerate([IMG_AHT20, IMG_LORA32, IMG_SETUP]):
        lx = Emu(228600) + (pw+gap)*i
        sl.shapes.add_picture(img, lx, pt, pw, ph)

    # Spec cards below photos
    cap_t = pt + ph + Emu(80000)
    cap_h = Emu(1500000)
    specs = [
        ("AHT20",
         "Interface: I²C (SoftI²C)\n"
         "Address:  0x38\n"
         "Pins: SDA=GPIO1  SCL=GPIO40\n"
         "Range: -40–85°C  0–100% RH\n"
         "Accuracy: ±0.3°C  ±2% RH"),
        ("Heltec WiFi LoRa32 V3",
         "MCU: ESP32-S3 @ 240 MHz, 8MB\n"
         "Radio: SX1262 LoRa 868 MHz +22dBm\n"
         "Display: SSD1306 OLED 128×64 px\n"
         "Serial: CP2102 USB-C → Pi\n"
         "FW: MicroPython 1.23"),
        ("Raspberry Pi 4B",
         "OS: Raspberry Pi OS 64-bit\n"
         "lora_bridge.py — reads LoRa32\n"
         "sentinel_client.py — push metrics\n"
         "Receives recovery cmds from hub\n"
         "Linked via Wi-Fi hotspot to Mac"),
    ]
    for i, (lbl, body) in enumerate(specs):
        lx = Emu(228600) + (pw+gap)*i
        _rect(sl, lx, cap_t, pw, cap_h, C_CARD)
        _tb(sl, lx+Emu(91440), cap_t+Emu(73152), pw-Emu(182880), Emu(246221),
            lbl, sz_emu=127000, bold=True, color=C_BLUE)
        _tb(sl, lx+Emu(91440), cap_t+Emu(360000), pw-Emu(182880), cap_h-Emu(420000),
            body, sz_emu=107950, color=C_DARK)


def build_hw_constraints(prs):
    """Hardware Constraints — 4 cards row + 2 cards row below."""
    sl = _blank(prs)
    _header(sl, "Hardware Constraints")
    _footer(sl, "Physical limits of embedded hardware shaped every design decision in the project")

    top1 = Emu(1180000)
    h1   = Emu(2000000)
    _four_cards(sl, top1, h1, [
        ("Power Budget (LoRa32)",
         "5 V / 500 mA via USB-C\n"
         "Peak under stress: ~350 mA\n"
         "LoRa TX burst: +22 dBm = 140 mA\n"
         "Solution: duty-cycle radio,\n"
         "deep-sleep between readings\n"
         "INA219 monitors voltage sag"),
        ("Processing (ESP32-S3)",
         "Dual-core 240 MHz Xtensa LX7\n"
         "RAM: 512 KB internal SRAM\n"
         "No OS scheduler — bare loop\n"
         "No ML on-device possible\n"
         "All inference delegated to hub\n"
         "Stress test proves recovery works"),
        ("I/O Pin Conflicts",
         "SX1262 LoRa: GPIO 8/9/10/12/14\n"
         "AHT20 SoftI²C: GPIO 1 & 40\n"
         "OLED HW I²C: GPIO 17, 18, 21\n"
         "USB UART0: GPIO 43/44\n"
         "All buses verified non-conflicting\n"
         "via I²C scan on every boot"),
        ("Radio Constraints (SX1262)",
         "868 MHz, BW=125 kHz, SF7, CR4/5\n"
         "Max payload: 222 bytes\n"
         "EU868 duty cycle: 1% max\n"
         "Current transport: USB serial\n"
         "(no duty-cycle limit)\n"
         "LoRa reserved for mesh future"),
    ])

    top2 = top1 + h1 + Emu(80000)
    h2   = Emu(1800000)

    # 2 wider cards for sensor accuracy + wiring
    cw = (Emu(12188952) - Emu(228600)*2 - Emu(91440)) // 2
    gap= Emu(91440)
    l0 = Emu(228600)

    for i, (lbl, body) in enumerate([
        ("Sensor Accuracy Limits",
         "AHT20: ±0.3 °C accuracy — adequate for crop/industrial use\n"
         "AHT20: ±2% RH — adequate for humidity monitoring\n"
         "No ADC on LoRa32 for analog sensors — I²C only\n"
         "Voltage monitoring: simulated on dev; real INA219 on Pi\n"
         "All sensor readings sanity-checked in firmware before pushing to hub"),
        ("Physical Wiring & PCB Limits",
         "Dev board — no custom PCB designed for this iteration\n"
         "Breadboard-style jumper wires: risk of loose contacts\n"
         "Resolved by using solder-tipped jumper connectors\n"
         "OLED RST pin requires explicit GPIO low→high pulse on boot\n"
         "USB re-enumeration required after each power cycle of LoRa32"),
    ]):
        lx = l0 + (cw+gap)*i
        _rect(sl, lx, top2, cw, h2, C_CARD)
        _tb(sl, lx+Emu(91440), top2+Emu(73152), cw-Emu(182880), Emu(246221),
            lbl, sz_emu=127000, bold=True, color=C_BLUE)
        _tb(sl, lx+Emu(91440), top2+Emu(360000), cw-Emu(182880), h2-Emu(420000),
            body, sz_emu=107950, color=C_DARK)


def build_sw_constraints(prs):
    """Software Constraints — 4 cards row + 2 cards below."""
    sl = _blank(prs)
    _header(sl, "Software Constraints")
    _footer(sl, "All constraints resolved in the current implementation — system is stable under continuous load")

    top1 = Emu(1180000)
    h1   = Emu(2000000)
    _four_cards(sl, top1, h1, [
        ("Real-Time Data Collection",
         "Metric loop must complete < 5 s\n"
         "psutil calls are blocking → run in thread\n"
         "In-memory event bus (no Redis needed)\n"
         "Serial read: 2 s timeout, auto-reconnect\n"
         "AI diagnosis in background thread\n"
         "— never blocks the metric pipeline"),
        ("False-Positive Prevention",
         "Warmup gate: suppress first 75 s\n"
         "Min 2 consecutive readings before alert\n"
         "1-min cooldown per metric after firing\n"
         "Baseline freeze during active anomaly\n"
         "Hysteresis: reset only at mean+0.5σ\n"
         "5 detection methods cross-checked"),
        ("LLM Rate Limits (Groq)",
         "Free tier: 30 req/min, 6000 tok/min\n"
         "Diagnosis fires max 1× per anomaly\n"
         "Fallback chain: Groq → Ollama → rules\n"
         "Ollama: local LLM, no rate limit\n"
         "Timeout 30 s → falls back to rule-based\n"
         "Never blocks event processing"),
        ("Multi-Device Coordination",
         "Each device has its own anomaly baseline\n"
         "Per-device cooldown & escalation state\n"
         "Hub serialises all events on one bus\n"
         "Device heartbeat timeout: 60 s\n"
         "Commands queued if device offline\n"
         "cmd_port=0 forces queue-only delivery"),
    ])

    top2 = top1 + h1 + Emu(80000)
    h2   = Emu(1800000)
    cw = (Emu(12188952) - Emu(228600)*2 - Emu(91440)) // 2
    gap= Emu(91440)
    l0 = Emu(228600)

    for i, (lbl, body) in enumerate([
        ("Security Constraints",
         "Demo mode: no auth on local network (acceptable for prototype)\n"
         "Claude AI threat analysis requires ANTHROPIC_API_KEY env var\n"
         "Allowlist ports: 22, 80, 443, 1883, 5001, 8883\n"
         "Synthetic threats injected at 4% probability for demo visibility\n"
         "Production path: add JWT tokens, HTTPS, Suricata/Zeek IDS integration"),
        ("Persistence & Crash Safety",
         "SQLite WAL mode — no data corruption on unexpected shutdown\n"
         "Recovery cooldown 60 s (demo) / 300 s (production) via config\n"
         "Graduated escalation L1→L4 per metric category\n"
         "Outcome verification 30 s after each recovery action\n"
         "Escalation state resets automatically when metric recovers below threshold"),
    ]):
        lx = l0 + (cw+gap)*i
        _rect(sl, lx, top2, cw, h2, C_CARD)
        _tb(sl, lx+Emu(91440), top2+Emu(73152), cw-Emu(182880), Emu(246221),
            lbl, sz_emu=127000, bold=True, color=C_BLUE)
        _tb(sl, lx+Emu(91440), top2+Emu(360000), cw-Emu(182880), h2-Emu(420000),
            body, sz_emu=107950, color=C_DARK)


def build_timeline(prs):
    """Project Timeline — 4 phase column cards + stat row."""
    sl = _blank(prs)
    _header(sl, "Project Timeline — 44 Days  ·  58 Commits")
    _footer(sl, "First commit: 03 Mar 2026  ·  Final demo: 16 Apr 2026  ·  git log --oneline | wc -l = 58")

    # 4 phase cards
    gap = Emu(91440)
    cw  = (Emu(12188952) - Emu(228600)*2 - gap*3) // 4
    ch  = Emu(3600000)
    ct  = Emu(1143000)
    l0  = Emu(228600)

    phases = [
        ("Week 1–2\n03 – 14 Mar",
         "Foundation",
         "Project scaffold\nConfig YAML + event bus\n"
         "Monitoring agent (psutil)\nCPU / Memory / Disk /\nNetwork metrics every 5 s\n"
         "Flask dashboard skeleton\nSQLite incident database"),
        ("Week 3–4\n15 – 28 Mar",
         "Intelligence",
         "Adaptive anomaly detection\n(z-score, IQR, trend)\n"
         "Rule-based diagnosis YAML\nOllama LLM integration\n"
         "Recovery agent 15+ actions\nGraduated escalation L1→L4\n"
         "Outcome verification 30 s"),
        ("Week 5–6\n29 Mar – 11 Apr",
         "Hardware & Multi-Device",
         "Heltec LoRa32 MicroPython FW\n"
         "AHT20 I²C sensor driver\nSSD1306 OLED display\n"
         "lora_bridge.py serial bridge\nRaspberry Pi integration\n"
         "Distributed dashboard panel\nPower monitoring agent"),
        ("Week 7\n12 – 16 Apr",
         "Integration & Demo",
         "CPU spike via serial CMD\nFull pipeline on 3 devices\n"
         "Groq LLaMA-3.3-70b AI\n"
         "Security threat agent\nKeras LSTM autoencoder\n"
         "Demo recording\nFinal documentation"),
    ]

    for i, (period, phase, detail) in enumerate(phases):
        lx = l0 + (cw+gap)*i
        # Phase header
        _rect(sl, lx, ct, cw, Emu(240000), C_BLUE)
        _tb(sl, lx+Emu(45720), ct+Emu(27216), cw-Emu(91440), Emu(200000),
            period, sz_emu=101600, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Phase name
        _rect(sl, lx, ct+Emu(240000), cw, Emu(230000), C_NAV)
        _tb(sl, lx+Emu(45720), ct+Emu(263000), cw-Emu(91440), Emu(200000),
            phase, sz_emu=114300, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Detail body
        _rect(sl, lx, ct+Emu(470000), cw, ch-Emu(470000), C_CARD)
        _tb(sl, lx+Emu(91440), ct+Emu(530000), cw-Emu(182880), ch-Emu(600000),
            detail, sz_emu=107950, color=C_DARK)

    # Stats row (matches exact positions from slide 2)
    _stat_box(sl, 0, "58",   "git\ncommits")
    _stat_box(sl, 1, "44",   "days\ntotal")
    _stat_box(sl, 2, "3",    "live\ndevices")
    _stat_box(sl, 3, "6",    "AI\nagents")


def build_team(prs):
    """Team Member Work Contributions — table style matching slide 5."""
    sl = _blank(prs)
    _header(sl, "Team Member Work Contributions")
    _footer(sl, "Solo project — full-stack embedded systems, AI pipeline, hardware integration and demo")

    # Member banner
    _rect(sl, Emu(228600), Emu(1143000), Emu(11731752), Emu(280000), C_NAV)
    _tb(sl, Emu(320040), Emu(1170000), Emu(11550000), Emu(240000),
        "Karthick S   —   Embedded Systems & IoT  ·  April 2026",
        sz_emu=139700, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Table: Area | Tasks Done | Output
    c1l=Emu(228600); c1w=Emu(2600000)
    c2l=Emu(2868600); c2w=Emu(5400000)
    c3l=Emu(8308600); c3w=Emu(3651352)
    hdr_t = Emu(1480000)
    hdr_h = Emu(265608)

    for l, w, txt in [(c1l,c1w,"Area"),(c2l,c2w,"Work Done"),(c3l,c3w,"Output / Evidence")]:
        _rect(sl, l, hdr_t, w, hdr_h, C_BLUE)
        _tb(sl, l+Emu(45720), hdr_t+Emu(27216), w-Emu(91440), hdr_h-Emu(54432),
            txt, sz_emu=114300, bold=True, color=C_WHITE)

    rows = [
        ("Firmware",       "MicroPython on ESP32-S3; AHT20 I²C driver with soft-reset,\ncalibration check, auto-retry; SSD1306 OLED; serial CMD handler",        "main.py on LoRa32"),
        ("Hardware Bridge","lora_bridge.py: serial reader, JSON parser, device register,\ncmd_port=0 queue, command poller thread every 3 s",                     "lora_bridge.py"),
        ("AI Pipeline",    "6 agents: monitoring, anomaly (5 methods), diagnosis (Groq +\nOllama), recovery (15+ actions, L1-L4), learning, security",           "agents/ directory"),
        ("Dashboard",      "Flask :5001 with Chart.js live charts, distributed device panel,\nAHT20 sensor tile, toast notifications, simulation lab",           "dashboard/app.py"),
        ("Multi-Device",   "Hub API: /api/devices/register, /api/metrics/push,\n/api/devices/<id>/commands — all three devices connected live",                 "sentinel_client.py"),
        ("Testing & Demo", "CPU spike, memory pressure, disk fill, power sag simulations;\nfull pipeline verified: anomaly → diagnosis → recovery on 3 devices", "Screen recordings"),
    ]

    rh = Emu(TBL_ROW_H.emu + 100000) if hasattr(TBL_ROW_H, 'emu') else Emu(360000)
    rh = Emu(360000)
    for i, (a, b, c) in enumerate(rows):
        t = hdr_t + hdr_h + rh*i
        bg = C_CARD if i%2==0 else C_WHITE
        for l, w, txt in [(c1l,c1w,a),(c2l,c2w,b),(c3l,c3w,c)]:
            _rect(sl, l, t, w, rh, bg)
            _tb(sl, l+Emu(45720), t+Emu(27216), w-Emu(91440), rh-Emu(54432),
                txt, sz_emu=107950, color=C_DARK)


def build_future_work(prs):
    """Future Work — 2 columns of 5 cards each (matching card style)."""
    sl = _blank(prs)
    _header(sl, "Future Work")
    _footer(sl, "Current system is a working proof-of-concept — the production roadmap is clear")

    gap = Emu(91440)
    cw  = (Emu(12188952) - Emu(228600)*2 - gap) // 2
    ch  = Emu(960000)
    l0  = Emu(228600)
    ct  = Emu(1143000)

    left = [
        ("LoRa Mesh Transport",
         "Replace USB serial with SX1262 radio links. Multi-hop mesh across "
         "many nodes — no Pi bridge needed per device. True wireless IoT."),
        ("Over-the-Air (OTA) Firmware Updates",
         "Push new MicroPython firmware to LoRa32 over Wi-Fi or LoRa radio "
         "without physical USB access. Essential for field deployment."),
        ("Cloud Twin — AWS IoT Core",
         "Push telemetry to IoT Core → DynamoDB shadow documents. "
         "Enables remote config, persistent device state and CloudWatch dashboards."),
        ("LSTM Inference On-Device",
         "Quantise LSTM to INT8 TFLite Micro. Run time-series anomaly detection "
         "directly on ESP32-S3 — true edge AI without hub dependency."),
        ("Mobile Alert App",
         "React Native app subscribes to hub WebSocket. Push notifications via "
         "FCM/APNs for critical anomalies — instant alerts on any phone."),
    ]
    right = [
        ("Custom PCB Design",
         "Replace breadboard with a 2-layer PCB integrating AHT20, OLED, "
         "INA219 and voltage divider on one board for a production-ready node."),
        ("Suricata / Zeek IDS Integration",
         "Connect security agent to real IDS for actual network traffic analysis. "
         "Replace synthetic demo threats with live threat intelligence."),
        ("Kubernetes & Horizontal Scaling",
         "Containerise hub agents with Docker. Helm chart for one-command deploy. "
         "Scale to 100+ devices with per-device agent isolation."),
        ("Federated Learning",
         "Devices train local Isolation Forest models and share only weights. "
         "Privacy-preserving collaborative anomaly learning across the fleet."),
        ("Energy Harvesting",
         "Solar panel + supercapacitor for LoRa32. Wake-on-interrupt sleep mode. "
         "Target < 10 mW average — years of operation on a single charge."),
    ]

    for i, (lbl, body) in enumerate(left):
        t = ct + (ch+gap)*i
        _rect(sl, l0, t, cw, ch, C_CARD)
        _rect(sl, l0, t, Emu(55000), ch, C_BLUE)
        _tb(sl, l0+Emu(110000), t+Emu(60000), cw-Emu(160000), Emu(220000),
            lbl, sz_emu=114300, bold=True, color=C_BLUE)
        _tb(sl, l0+Emu(110000), t+Emu(295000), cw-Emu(160000), ch-Emu(350000),
            body, sz_emu=101600, color=C_DARK)

    for i, (lbl, body) in enumerate(right):
        t = ct + (ch+gap)*i
        rx = l0+cw+gap
        _rect(sl, rx, t, cw, ch, C_CARD)
        _rect(sl, rx, t, Emu(55000), ch, C_GREEN)
        _tb(sl, rx+Emu(110000), t+Emu(60000), cw-Emu(160000), Emu(220000),
            lbl, sz_emu=114300, bold=True, color=C_BLUE)
        _tb(sl, rx+Emu(110000), t+Emu(295000), cw-Emu(160000), ch-Emu(350000),
            body, sz_emu=101600, color=C_DARK)


# ── Reorder helper ────────────────────────────────────────────────────────────
def reorder(prs, order):
    xml = prs.slides._sldIdLst
    ids = list(xml)
    for el in ids:
        xml.remove(el)
    for i in order:
        xml.append(ids[i])


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    src = '/Users/karthi/Downloads/Sentinel_AI_Presentation  -  Repaired.pptx'
    prs = Presentation(src)

    # Existing indices:
    # 0=Title  1=ProjectOverview  2=HardwareSetup  3=DataSources
    # 4=Challenges  5=Thresholds  6=Monitors  7=AIAgents
    # 8=Testing  9=AIModels  10=Demo  11=References

    build_purpose(prs)         # → 12
    build_hardware_photos(prs) # → 13
    build_timeline(prs)        # → 14
    build_team(prs)            # → 15
    build_hw_constraints(prs)  # → 16
    build_sw_constraints(prs)  # → 17
    build_future_work(prs)     # → 18

    # Final slide order:
    # 1. Title
    # 2. Purpose
    # 3. Hardware Photos
    # 4. Project Overview (existing)
    # 5. Hardware & System Setup (existing)
    # 6. Hardware Constraints (new)
    # 7. Software Constraints (new)
    # 8. Where Does Data Come From (existing)
    # 9. How Thresholds Work (existing)
    # 10. How Monitors Work (existing)
    # 11. Role of AI Agents (existing)
    # 12. What AI Models (existing)
    # 13. How Are We Testing (existing)
    # 14. Key Challenges (existing)
    # 15. Project Timeline (new)
    # 16. Team Contributions (new)
    # 17. Future Work (new)
    # 18. Demo
    # 19. References
    reorder(prs, [0, 12, 13, 1, 2, 16, 17, 3, 5, 6, 7, 9, 8, 4, 14, 15, 18, 10, 11])

    out = '/Users/karthi/Desktop/Sentinal_AI/Sentinel_AI_Final.pptx'
    prs.save(out)

    # Print final order
    prs2 = Presentation(out)
    print(f"Saved: {out}  ({len(prs2.slides)} slides)\n")
    for i, sl in enumerate(prs2.slides):
        for sh in sl.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.paragraphs[0].text.strip()
                if t:
                    print(f"  {i+1:2d}. {t[:75]}")
                    break
        else:
            print(f"  {i+1:2d}. (title / image slide)")


if __name__ == '__main__':
    main()

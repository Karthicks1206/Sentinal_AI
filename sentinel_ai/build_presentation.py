#!/usr/bin/env python3
"""
Build Sentinel AI presentation PPTX from the test.pptx template design.
Usage: python build_presentation.py
Output: /Users/karthi/Desktop/Sentinal_AI/Sentinel_AI_Presentation.pptx
"""

import copy
import shutil
from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# ── Colour palette (extracted from test.pptx) ────────────────────────────────
C_NAV   = RGBColor(0x1E, 0x2A, 0x3A)   # header / footer bar
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE  = RGBColor(0x41, 0x6E, 0x8F)   # section labels, stat numbers
C_DARK  = RGBColor(0x2C, 0x3E, 0x50)   # body text
C_GREEN = RGBColor(0x1A, 0x9E, 0x8A)   # positive accent
C_RED   = RGBColor(0xC0, 0x39, 0x2B)   # negative accent
C_GRAY  = RGBColor(0x7F, 0x8C, 0x8D)   # footer text
C_CARD  = RGBColor(0xF4, 0xF6, 0xF8)   # card background

# ── Layout constants (EMU) ────────────────────────────────────────────────────
SW   = Emu(12188952)   # slide width
SH   = Emu(6858000)    # slide height
HDR  = Emu(1005840)    # header height
FTR  = Emu(6492240)    # footer top
FH   = Emu(365760)     # footer height
ML   = Emu(274320)     # left margin
CT   = Emu(1097280)    # content top
CW   = Emu(11640312)   # content width (SW - 2*ML)
CH   = Emu(5394960)    # content height (FTR - CT)


# ── Low-level helpers ─────────────────────────────────────────────────────────
def _add_rect(slide, left, top, width, height, fill_rgb, line=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not line:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text,
                 font_size_pt=9, bold=False, color=None, align=PP_ALIGN.LEFT,
                 wrap=True, word_wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


def _add_header(slide, title_text):
    """Standard dark-navy header bar with white title."""
    _add_rect(slide, Emu(0), Emu(0), SW, HDR, C_NAV)
    _add_textbox(slide, ML, Emu(137160), Emu(11430000), Emu(492443),
                 title_text, font_size_pt=26, bold=True, color=C_WHITE)


def _add_footer(slide, ref_text=""):
    """Standard footer bar."""
    _add_rect(slide, Emu(0), FTR, SW, FH, C_NAV)
    if ref_text:
        _add_textbox(slide, Emu(182880), Emu(6510528), Emu(11887200), Emu(320040),
                     ref_text, font_size_pt=8, color=C_GRAY)


def _add_card(slide, left, top, width, height, label, body, label_size=10, body_size=9):
    """Card: light-gray rectangle + label (blue) + body (dark)."""
    _add_rect(slide, left, top, width, height, C_CARD)
    _add_textbox(slide, left + Emu(91440), top + Emu(73152),
                 width - Emu(182880), Emu(246221),
                 label, font_size_pt=label_size, bold=True, color=C_BLUE)
    _add_textbox(slide, left + Emu(91440), top + Emu(365760),
                 width - Emu(182880), height - Emu(457200),
                 body, font_size_pt=body_size, color=C_DARK)


def _stat_card(slide, left, top, width, height, stat, label):
    """Stat card: big coloured number + small label."""
    _add_rect(slide, left, top, width, height, C_CARD)
    _add_textbox(slide, left + Emu(91440), top + Emu(73152),
                 width - Emu(182880), Emu(502920),
                 stat, font_size_pt=22, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, left + Emu(91440), top + Emu(640080),
                 width - Emu(182880), Emu(500000),
                 label, font_size_pt=9, color=C_DARK, align=PP_ALIGN.CENTER)


def _row_header(slide, cols, col_ws, top, row_h=Emu(320040)):
    """Table header row: solid C_NAV cells with white text."""
    x = ML
    for label, w in zip(cols, col_ws):
        _add_rect(slide, x, top, w, row_h, C_BLUE)
        _add_textbox(slide, x + Emu(45720), top + Emu(27216),
                     w - Emu(91440), row_h - Emu(54432),
                     label, font_size_pt=9, bold=True, color=C_WHITE)
        x += w


def _table_row(slide, cells, col_ws, top, row_h=Emu(338328), bg=None):
    """Table data row."""
    x = ML
    bg_color = bg or C_CARD
    for text, w in zip(cells, col_ws):
        _add_rect(slide, x, top, w, row_h, bg_color)
        _add_textbox(slide, x + Emu(45720), top + Emu(36576),
                     w - Emu(91440), row_h - Emu(73152),
                     text, font_size_pt=8.5, color=C_DARK)
        x += w


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide — copied from template slide 1 (keeps background image)."""
    template = Presentation('/Users/karthi/Downloads/test.pptx')
    src_slide = template.slides[0]
    # Deep-copy the slide XML into a new blank slide
    blank_layout = prs.slide_layouts[6]  # Blank
    new_slide = prs.slides.add_slide(blank_layout)
    # Remove default shapes
    sp_tree = new_slide.shapes._spTree
    for ch in list(sp_tree):
        sp_tree.remove(ch)
    # Copy all elements from source slide
    for el in src_slide.shapes._spTree:
        sp_tree.append(copy.deepcopy(el))
    # Copy relationships (images etc.)
    for rel in src_slide.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_part = rel.target_part
                new_rel = new_slide.part.relate_to(img_part, rel.reltype)
                # Fix rId references in copied XML
                old_rid = rel.rId
                for el in sp_tree.iter():
                    for attr in list(el.attrib):
                        if el.attrib[attr] == old_rid:
                            el.attrib[attr] = new_rel
            except Exception:
                pass
    # Add title overlay text
    _add_textbox(new_slide, Emu(500000), Emu(2200000), Emu(7500000), Emu(800000),
                 "Sentinel AI",
                 font_size_pt=40, bold=True, color=C_WHITE)
    _add_textbox(new_slide, Emu(500000), Emu(3050000), Emu(7500000), Emu(500000),
                 "Self-Healing Distributed IoT Infrastructure",
                 font_size_pt=20, bold=False, color=RGBColor(0xA8, 0xC8, 0xE8))
    _add_textbox(new_slide, Emu(500000), Emu(3600000), Emu(7500000), Emu(400000),
                 "Embedded Systems & IoT  |  April 2026  |  Karthick S",
                 font_size_pt=14, color=C_GRAY)


def slide_purpose(prs):
    """Purpose & Motivation — agriculture IoT framing, limitless applications."""
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Why Does This System Exist?")
    _add_footer(sl, "Temperature & humidity are critical control variables across nearly every industry")

    # Big problem statement
    _add_rect(sl, ML, CT, CW, Emu(310000), RGBColor(0x41, 0x6E, 0x8F))
    _add_textbox(sl, ML + Emu(91440), CT + Emu(55000), CW - Emu(182880), Emu(220000),
                 "IoT devices fail silently. By the time an operator notices, the damage is already done.",
                 font_size_pt=13, bold=True, color=C_WHITE)

    gap = Emu(91440)

    # Left column — Agriculture use case (primary)
    col_w = Emu(5600000)
    col_h = Emu(4400000)
    _add_rect(sl, ML, CT + Emu(380000), col_w, col_h, C_CARD)

    # Green accent bar on left edge
    _add_rect(sl, ML, CT + Emu(380000), Emu(55000), col_h, C_GREEN)

    _add_textbox(sl, ML + Emu(150000), CT + Emu(450000), col_w - Emu(200000), Emu(280000),
                 "Primary Application — Precision Agriculture",
                 font_size_pt=13, bold=True, color=C_GREEN)

    agri_text = (
        "Crops are highly sensitive to temperature and humidity. A cold snap,\n"
        "unexpected heat spike, or humidity surge can destroy an entire harvest.\n\n"
        "The AHT20 sensor monitors both continuously:\n"
        "   Temp °C  — frost risk, heat stress, cold-chain compliance\n"
        "   Humidity % — mould growth, irrigation need, VPD control\n\n"
        "Sentinel AI watches these readings in real-time. When a value\n"
        "drifts outside the learned normal range:\n\n"
        "   1.  DETECT   — anomaly fires within 10 seconds\n"
        "   2.  DIAGNOSE — AI identifies root cause\n"
        "   3.  ALERT    — toast notification to operator\n"
        "   4.  RECOVER  — automated corrective action\n"
        "   5.  LOG      — incident stored for pattern learning\n\n"
        "No human needed to be watching. The system heals itself."
    )
    _add_textbox(sl, ML + Emu(150000), CT + Emu(780000),
                 col_w - Emu(200000), col_h - Emu(450000),
                 agri_text, font_size_pt=9, color=C_DARK)

    # Right column — Applications are limitless
    right_x = ML + col_w + gap
    right_w = CW - col_w - gap
    _add_rect(sl, right_x, CT + Emu(380000), right_w, col_h, C_CARD)
    _add_rect(sl, right_x, CT + Emu(380000), Emu(55000), col_h, C_BLUE)

    _add_textbox(sl, right_x + Emu(150000), CT + Emu(450000),
                 right_w - Emu(200000), Emu(280000),
                 "Applications Are Limitless",
                 font_size_pt=13, bold=True, color=C_BLUE)

    apps = [
        ("Agriculture",      "Greenhouse climate control, cold-chain,\ncrop disease prevention"),
        ("Industrial IoT",   "Factory floor heat monitoring, motor\ntemperature, predictive maintenance"),
        ("Healthcare",       "Drug cold-chain, sterile room humidity,\npatient environment monitoring"),
        ("Data Centres",     "Rack inlet temp, cooling efficiency,\nfire-risk humidity control"),
        ("Smart Buildings",  "HVAC optimisation, occupant comfort,\nenergy usage reduction"),
        ("Environmental",    "Wildlife habitat sensors, flood\nmonitoring, weather stations"),
    ]

    app_h = Emu(560000)
    for i, (title, desc) in enumerate(apps):
        top = CT + Emu(800000) + app_h * i
        _add_textbox(sl, right_x + Emu(150000), top,
                     right_w - Emu(200000), Emu(230000),
                     title, font_size_pt=9.5, bold=True, color=C_BLUE)
        _add_textbox(sl, right_x + Emu(150000), top + Emu(240000),
                     right_w - Emu(200000), Emu(300000),
                     desc, font_size_pt=8.5, color=C_DARK)


def slide_hardware_photos(prs):
    """Real photos of the hardware — AHT20, LoRa32, full setup."""
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "The Hardware — What We Actually Built")
    _add_footer(sl, "Left: AHT20 I²C sensor  ·  Centre: Heltec WiFi LoRa32 V3 (ESP32-S3)  ·  Right: Full setup — LoRa32 OLED showing 'Sentinel AI' + Raspberry Pi")

    IMG_AHT20  = "/tmp/hw_imgs/aht20.jpg"
    IMG_LORA32 = "/tmp/hw_imgs/lora32.jpg"
    IMG_SETUP  = "/tmp/hw_imgs/setup.jpg"

    gap = Emu(91440)
    # Three equal photo columns
    photo_w = (CW - gap * 2) // 3
    photo_h = Emu(3800000)
    photo_top = CT + Emu(80000)

    sl.shapes.add_picture(IMG_AHT20,  ML,                     photo_top, photo_w, photo_h)
    sl.shapes.add_picture(IMG_LORA32, ML + photo_w + gap,     photo_top, photo_w, photo_h)
    sl.shapes.add_picture(IMG_SETUP,  ML + (photo_w + gap)*2, photo_top, photo_w, photo_h)

    # Caption cards under each photo
    cap_h = Emu(1200000)
    cap_top = photo_top + photo_h + Emu(60000)

    _add_card(sl, ML, cap_top, photo_w, cap_h,
              "AHT20 — Temp & Humidity Sensor",
              "Adafruit breakout board\n"
              "I²C interface (4-wire: VCC / GND / SDA / SCL)\n"
              "Address: 0x38\n"
              "Accuracy: ±0.3 °C   ±2% RH\n"
              "Connected to LoRa32 GPIO 1 & 40",
              body_size=8.5)

    _add_card(sl, ML + photo_w + gap, cap_top, photo_w, cap_h,
              "Heltec WiFi LoRa32 V3",
              "ESP32-S3 @ 240 MHz · 8 MB flash\n"
              "On-board SX1262 LoRa radio (868 MHz)\n"
              "Built-in SSD1306 OLED 128×64 px\n"
              "USB-C serial (CP2102) → Pi bridge\n"
              "Running MicroPython 1.23",
              body_size=8.5)

    _add_card(sl, ML + (photo_w + gap)*2, cap_top, photo_w, cap_h,
              "Full Sentinel AI Node Setup",
              "LoRa32 OLED shows: Sentinel AI\n"
              "CPU / Mem / Temp / Hum live\n"
              "USB-C cable → Raspberry Pi 4B\n"
              "Pi runs lora_bridge.py + sentinel_client.py\n"
              "Pi visible top-right with USB hub",
              body_size=8.5)


def slide_02_iot_overview(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "IoT Device Overview — The Full System")
    _add_footer(sl, "End-to-end self-healing pipeline across heterogeneous hardware")

    # Pipeline label
    _add_textbox(sl, ML, CT + Emu(0), CW, Emu(274320),
                 "Sense  →  Transmit  →  Detect  →  Diagnose  →  Recover",
                 font_size_pt=13, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    gap = Emu(91440)
    card_w = (CW - gap * 2) // 3
    card_h = Emu(1500000)
    top = CT + Emu(310000)

    _add_card(sl, ML, top, card_w, card_h,
              "LoRa32 Node (ESP32-S3)",
              "MicroPython firmware\n"
              "AHT20 sensor: T°C + RH%\n"
              "SSD1306 OLED display\n"
              "USB serial → Pi bridge\n"
              "CPU/mem simulated stress",
              label_size=11, body_size=9)

    _add_card(sl, ML + card_w + gap, top, card_w, card_h,
              "Raspberry Pi 4 (Hub Node)",
              "Linux · Python 3\n"
              "lora_bridge.py serial reader\n"
              "sentinel_client.py metric push\n"
              "Forwards metrics via HTTP POST\n"
              "Receives recovery commands",
              label_size=11, body_size=9)

    _add_card(sl, ML + (card_w + gap) * 2, top, card_w, card_h,
              "Mac (AI Engine / Hub)",
              "main.py — 6 AI agents\n"
              "Flask dashboard :5001\n"
              "Ollama LLaMA-3.2:3b (local)\n"
              "Groq LLaMA-3.3-70b (cloud)\n"
              "SQLite incident log",
              label_size=11, body_size=9)

    # Stats row
    stat_w = CW // 5
    stats_top = top + card_h + Emu(120000)
    for i, (s, l) in enumerate([("3", "physical\ndevices"), ("6", "AI\nagents"),
                                  ("5", "detection\nmethods"), ("<30s", "detect\n→ fix"),
                                  ("15+", "recovery\nactions")]):
        _stat_card(sl, ML + stat_w * i, stats_top, stat_w - Emu(45720), Emu(900000), s, l)


def slide_03_hardware_summary(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Hardware Summary")
    _add_footer(sl, "All physical hardware integrated and tested 09 Apr 2026")

    cols = ["Component", "Role", "Interface", "Key Spec / Notes"]
    col_ws = [Emu(2400000), Emu(2800000), Emu(2000000), Emu(4440312)]

    _row_header(sl, cols, col_ws, CT)

    rows = [
        ("Heltec WiFi LoRa32 V3",   "Edge node — sense + transmit",      "USB-C serial (CP2102)", "ESP32-S3, 240 MHz, 8 MB flash, SX1262 LoRa"),
        ("Raspberry Pi 4B",          "Bridge + remote metric node",        "Ethernet / Wi-Fi",      "ARM Cortex-A72 64-bit, 4 GB RAM, Linux"),
        ("AHT20",                    "Temperature & humidity sensor",      "I²C (SDA=1, SCL=40)",   "0x38, ±0.3°C / ±2% RH accuracy"),
        ("SSD1306 OLED",             "Local display on LoRa32",            "I²C (SDA=17, SCL=18)",  "128×64 px, RST=GPIO21, 3.3 V"),
        ("SX1262 (on LoRa32 V3)",    "Long-range radio (future transport)","SPI (internal)",        "868 MHz, +22 dBm, SF7–SF12"),
        ("INA219 (simulated on Pi)", "Power monitoring",                   "I²C",                   "Voltage / current / watt readings"),
        ("CP2102 USB-UART",          "Serial bridge chip",                 "USB-C to Pi USB-A",     "115200 baud, auto-detected by lora_bridge.py"),
        ("MacBook Pro M3",           "AI hub & dashboard server",          "Wi-Fi / USB",           "Flask :5001, Ollama, SQLite, Chart.js"),
    ]

    rh = Emu(338328)
    for i, row in enumerate(rows):
        bg = C_CARD if i % 2 == 0 else C_WHITE
        _table_row(sl, row, col_ws, CT + Emu(320040) + rh * i, rh, bg)


def slide_04_timeline(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Project Timeline — 44 Days · 58 Commits")
    _add_footer(sl, "git log --oneline | wc -l = 58 commits  ·  first commit 03 Mar 2026  ·  demo 16 Apr 2026")

    phases = [
        ("Week 1–2\n03–14 Mar", "Foundation",
         "Project scaffold · config YAML · event bus · monitoring agent\n"
         "CPU / memory / disk / network metrics via psutil"),
        ("Week 3–4\n15–28 Mar", "Intelligence",
         "Anomaly detection (z-score, IQR, trend) · rule-based diagnosis\n"
         "Ollama LLM integration · recovery agent (15+ actions)"),
        ("Week 5–6\n29 Mar–11 Apr", "Hardware & Multi-device",
         "Heltec LoRa32 MicroPython firmware · AHT20 sensor · OLED display\n"
         "Raspberry Pi bridge · distributed dashboard panel"),
        ("Week 7\n12–16 Apr", "Integration & Demo",
         "CPU spike via serial CMD · full pipeline verified on 3 devices\n"
         "Groq LLaMA-3.3-70b · security agent · demo recording"),
    ]

    card_h = Emu(1250000)
    card_w = (CW - Emu(91440) * 3) // 4
    gap = Emu(91440)

    for i, (period, phase, detail) in enumerate(phases):
        left = ML + (card_w + gap) * i
        _add_rect(sl, left, CT, card_w, Emu(250000),
                  RGBColor(0x41, 0x6E, 0x8F))
        _add_textbox(sl, left + Emu(45720), CT + Emu(27216),
                     card_w - Emu(91440), Emu(220000),
                     period, font_size_pt=8.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _add_rect(sl, left, CT + Emu(250000), card_w, card_h, C_CARD)
        _add_textbox(sl, left + Emu(45720), CT + Emu(275000),
                     card_w - Emu(91440), Emu(230000),
                     phase, font_size_pt=11, bold=True, color=C_BLUE)
        _add_textbox(sl, left + Emu(45720), CT + Emu(530000),
                     card_w - Emu(91440), card_h - Emu(320000),
                     detail, font_size_pt=8.5, color=C_DARK)

    # Stat strip
    stats_top = CT + Emu(1570000)
    stat_w = CW // 4
    for i, (s, l) in enumerate([("58", "git commits"), ("44", "days total"),
                                  ("3", "live devices"), ("6", "AI agents")]):
        _stat_card(sl, ML + stat_w * i, stats_top, stat_w - Emu(45720), Emu(900000), s, l)

    # Contributions
    contrib_top = stats_top + Emu(980000)
    _add_textbox(sl, ML, contrib_top, CW, Emu(250000),
                 "Team Contribution", font_size_pt=11, bold=True, color=C_BLUE)
    _add_textbox(sl, ML, contrib_top + Emu(265000), CW // 2, Emu(700000),
                 "Karthick S — Full-stack development: firmware, bridge, AI pipeline,\n"
                 "dashboard, multi-device coordination, testing & demo",
                 font_size_pt=9, color=C_DARK)


def slide_05_sensor_interfacing(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Sensor Hardware & Interfacing")
    _add_footer(sl, "I²C protocol: 7-bit addressing, 100 kHz SCL, ACK/NACK, repeated-start for reads")

    gap = Emu(91440)
    card_w = (CW - gap) // 2
    card_h = Emu(2000000)

    # AHT20 card
    _add_card(sl, ML, CT, card_w, card_h,
              "AHT20 — Temperature & Humidity",
              "Protocol:  I²C  (SoftI²C on ESP32-S3)\n"
              "Address:   0x38\n"
              "SDA pin:   GPIO 1      SCL pin: GPIO 40\n"
              "Frequency: 100 kHz\n\n"
              "Init sequence:\n"
              "  1. Soft reset  (0xBA)\n"
              "  2. Read status byte — check CAL bit (0x08)\n"
              "  3. Send calibrate cmd (0xBE 0x08 0x00) if needed\n\n"
              "Measurement:\n"
              "  1. Trigger measurement (0xAC 0x33 0x00)\n"
              "  2. Wait 80 ms\n"
              "  3. Read 7 bytes → 20-bit T + 20-bit RH\n"
              "  4. Sanity-check: 0°C–60°C, 0%–100% RH\n\n"
              "Output: T = 25.5°C  ·  RH = 55%",
              label_size=11, body_size=8.5)

    # OLED card
    _add_card(sl, ML + card_w + gap, CT, card_w, card_h,
              "SSD1306 OLED — 128×64 px Local Display",
              "Protocol:  I²C\n"
              "Address:   0x3C\n"
              "SDA pin:   GPIO 17     SCL pin: GPIO 18\n"
              "RST pin:   GPIO 21\n"
              "Supply:    3.3 V\n\n"
              "Display content (updated every 2 s):\n"
              "  Line 1: CPU % + memory %\n"
              "  Line 2: Temperature °C + humidity %\n"
              "  Line 3: Status (OK / ANOMALY / RECOVERY)\n\n"
              "Driver: MicroPython ssd1306 framebuf library\n"
              "Init: reset pulse → send init cmd sequence",
              label_size=11, body_size=8.5)

    # Protocol detail strip
    strip_top = CT + card_h + Emu(120000)
    strip_h = Emu(1500000)
    sub_w = CW // 3

    _add_card(sl, ML, strip_top, sub_w - gap, strip_h,
              "I²C Bus Timing",
              "Start: SDA falls while SCL high\n"
              "Data: SDA stable while SCL high\n"
              "ACK: receiver pulls SDA low\n"
              "Stop: SDA rises while SCL high",
              body_size=8.5)

    _add_card(sl, ML + sub_w, strip_top, sub_w - gap, strip_h,
              "Error Handling",
              "I²C scan on boot → log all\nfound addresses\n"
              "Auto-retry 3× if sensor\nnot found at boot\n"
              "Fallback: display '—' for\nmissing readings",
              body_size=8.5)

    _add_card(sl, ML + sub_w * 2, strip_top, sub_w - gap, strip_h,
              "Power Monitoring (INA219)",
              "Simulated on dev hardware\n"
              "Real hardware: INA219 I²C\n"
              "Reads: voltage V, current A,\npower W every 5 s\n"
              "Alert: ±10% of 5 V nominal",
              body_size=8.5)


def slide_06_display_output(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Display Output Approach")
    _add_footer(sl, "Three complementary output layers — hardware, web, structured log")

    gap = Emu(91440)
    card_w = (CW - gap * 2) // 3
    card_h = Emu(2800000)

    _add_card(sl, ML, CT, card_w, card_h,
              "OLED (SSD1306) — Edge Node",
              "128×64 px monochrome display\n"
              "on the LoRa32 itself\n\n"
              "Updates every 2 seconds:\n"
              "  • CPU %   Memory %\n"
              "  • Temp °C   Humidity %\n"
              "  • Status line\n\n"
              "No network needed — always\nshows live local readings\n\n"
              "Uses MicroPython\nframebuf + ssd1306 driver",
              label_size=11, body_size=9)

    _add_card(sl, ML + card_w + gap, CT, card_w, card_h,
              "Flask Web Dashboard — Hub",
              "Browser at http://hub:5001\n\n"
              "Live panels:\n"
              "  • CPU / Memory / Disk / Network\n"
              "  • Power voltage + quality chart\n"
              "  • Distributed devices tab\n"
              "  • AHT20 sensor tile (orange)\n"
              "  • Incident timeline\n"
              "  • AI agent status\n\n"
              "Chart.js auto-refresh every 5 s\n"
              "Toast notifications (no-click)\n"
              "Simulation Lab buttons",
              label_size=11, body_size=9)

    _add_card(sl, ML + (card_w + gap) * 2, CT, card_w, card_h,
              "Structured JSON Logs",
              "logs/sentinel.log\n"
              "Format: JSON-lines\n\n"
              "Fields per event:\n"
              "  • timestamp (ISO-8601 UTC)\n"
              "  • level (INFO / WARNING)\n"
              "  • device_id\n"
              "  • event type (anomaly /\n"
              "    diagnosis / recovery)\n"
              "  • metric values\n\n"
              "Retention: 100 MB rolling\n"
              "10 backup files\n"
              "SQLite incident DB",
              label_size=11, body_size=9)


def slide_07_stack(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "OS, Languages & Frameworks — Design Choices")
    _add_footer(sl, "Chosen for reliability, zero licensing cost, and direct hardware access")

    cols = ["Layer", "Choice", "Why We Chose It"]
    col_ws = [Emu(2200000), Emu(2600000), Emu(6840312)]
    _row_header(sl, cols, col_ws, CT)

    rows = [
        ("Edge OS",          "MicroPython 1.23 (ESP32-S3)",
         "Runs directly on bare metal — no RTOS needed. Built-in machine, SoftI2C, "
         "ssd1306 libs. OTA-updatable over serial. Tiny footprint (≈180 KB)."),
        ("Hub OS",           "Linux (Raspberry Pi OS) + macOS",
         "POSIX signals, cron, systemctl. Full Python 3.12 ecosystem. "
         "psutil gives CPU/mem/disk/net without root."),
        ("Firmware lang",    "MicroPython",
         "Python syntax on microcontrollers. GPIO, I²C, SPI via machine module. "
         "Interactive REPL over serial for live debugging."),
        ("Backend lang",     "Python 3.12",
         "asyncio + threads. Rich ML libs (sklearn, keras, torch). "
         "Same language end-to-end — no context switch."),
        ("Web framework",    "Flask 3.0",
         "Lightweight, no magic. Single-file dashboard app. "
         "Server-Sent Events for live data push to browser."),
        ("ML / AI",          "sklearn · Keras · PyTorch (MPS) · Ollama · Groq",
         "Isolation Forest (multivariate) + LSTM autoencoder (time-series) on-device. "
         "Ollama: free local LLM. Groq: free-tier cloud LLM, 2–8 s response."),
        ("Database",         "SQLite 3 / PostgreSQL (configurable)",
         "Zero-config SQLite for dev/demo. Postgres for production. "
         "Schema shared — swap with one config flag."),
        ("Frontend",         "Vanilla JS + Chart.js 4",
         "No build step. Chart.js for live time-series. "
         "CSS-only toast notifications — no jQuery dependency."),
        ("Comms",            "HTTP REST (internal) · Serial UART · I²C",
         "REST: device → hub metric push and command delivery. "
         "UART 115200 baud: LoRa32 → Pi. I²C: sensors on board."),
    ]

    rh = Emu(330000)
    for i, row in enumerate(rows):
        bg = C_CARD if i % 2 == 0 else C_WHITE
        _table_row(sl, row, col_ws, CT + Emu(320040) + rh * i, rh, bg)


def slide_08_hw_constraints(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Hardware Constraints")
    _add_footer(sl, "Constraints drove design decisions throughout the project")

    gap = Emu(91440)
    card_w = (CW - gap) // 2
    card_h = Emu(1350000)

    cards = [
        ("Power Budget (LoRa32)",
         "5 V / 500 mA via USB-C\n"
         "Peak CPU stress: ~350 mA\n"
         "LoRa TX: +22 dBm = 140 mA burst\n"
         "Design: duty-cycle LoRa, idle deep-sleep\n"
         "Simulated INA219 monitors voltage sag"),
        ("Processing (ESP32-S3)",
         "Dual-core 240 MHz Xtensa\n"
         "RAM: 512 KB internal SRAM\n"
         "No OS scheduler — cooperative loop\n"
         "No ML on-device; all inference on hub\n"
         "Stress test proves recovery works"),
        ("I/O Pin Conflicts",
         "SX1262 LoRa uses GPIO 8/9/10/12/14\n"
         "AHT20 on GPIO 1/40 (SoftI²C)\n"
         "OLED on GPIO 17/18/21 (HW I²C)\n"
         "All buses verified non-conflicting\n"
         "USB serial occupies UART0"),
        ("Radio Constraints",
         "SX1262 @ 868 MHz, BW=125 kHz, SF7\n"
         "Payload limit: ~222 bytes\n"
         "Duty cycle: EU868 = 1% max\n"
         "Current transport: USB serial (no limit)\n"
         "LoRa reserved for future multi-node mesh"),
        ("Sensor Accuracy",
         "AHT20: ±0.3°C, ±2% RH\n"
         "No ADC on LoRa32 for analog sensors\n"
         "INA219 I²C for real power on Pi\n"
         "All readings sanity-checked in firmware\n"
         "Outliers discarded before push"),
        ("PCB / Wiring",
         "Dev board — no custom PCB\n"
         "Breadboard wiring risk: loose contacts\n"
         "Resolved: solder-tipped jumpers\n"
         "OLED RST required explicit GPIO pulse\n"
         "USB re-enumeration after power cycle"),
    ]

    for i, (label, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        left = ML + (card_w + gap) * col
        top = CT + (card_h + gap) * row
        _add_card(sl, left, top, card_w, card_h, label, body, body_size=8.5)


def slide_09_sw_constraints(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Software Constraints")
    _add_footer(sl, "All constraints resolved in current implementation")

    gap = Emu(91440)
    card_w = (CW - gap) // 2
    card_h = Emu(1350000)

    cards = [
        ("Real-time Collection",
         "Metric loop must complete in <5 s\n"
         "psutil calls are blocking — run in thread\n"
         "Event bus is in-memory (no Redis needed)\n"
         "Diagnosis in background thread (non-blocking)\n"
         "Serial read: 2 s timeout, reconnect on error"),
        ("False-positive Prevention",
         "Warmup gate: suppress anomalies first 75 s\n"
         "Min 2 consecutive readings before alert fires\n"
         "1-min cooldown per metric after firing\n"
         "Baseline freeze during active anomaly\n"
         "Hysteresis: reset only at mean+0.5σ"),
        ("LLM Rate Limits (Groq)",
         "Free tier: 30 req/min, 6000 tokens/min\n"
         "Diagnosis runs max 1× per anomaly\n"
         "Fallback chain: Groq → Ollama → rules\n"
         "Ollama: local LLM, no rate limit\n"
         "Timeout: 30 s — returns rule-based if hit"),
        ("Multi-device Coordination",
         "Each remote device gets its own baseline\n"
         "Per-device anomaly cooldown & escalation\n"
         "Hub serialises all events on one bus\n"
         "Device heartbeat timeout: 60 s\n"
         "Commands queued if device unreachable"),
        ("Security",
         "No auth on local dev (demo mode)\n"
         "Claude AI threat analysis via API key\n"
         "Allowlist ports: 22/80/443/1883/5001\n"
         "Synthetic threats injected for demo (4%)\n"
         "Production: add JWT + HTTPS + Suricata"),
        ("Persistence & Recovery",
         "SQLite WAL mode — no corruption on crash\n"
         "Recovery cooldown 60 s (demo) / 300 s (prod)\n"
         "Graduated escalation: L1→L4 per metric\n"
         "Outcome verification 30 s after action\n"
         "Escalation resets when metric recovers"),
    ]

    for i, (label, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        left = ML + (card_w + gap) * col
        top = CT + (card_h + gap) * row
        _add_card(sl, left, top, card_w, card_h, label, body, body_size=8.5)


def slide_10_challenges(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Key Challenges & Solutions")
    _add_footer(sl, "58 commits — most addressed at least one of the challenges below")

    cols = ["Challenge", "Root Cause", "Solution"]
    col_ws = [Emu(3200000), Emu(3600000), Emu(4840312)]
    _row_header(sl, cols, col_ws, CT)

    rows = [
        ("AHT20 showing T:-- H:--",
         "No calibration check, no retry on failed I²C scan",
         "Added soft-reset (0xBA), CAL-bit check, 3× auto-retry, I²C scan debug print"),
        ("CPU spike misrouted to Pi",
         "LoRa32 registered with cmd_port=5002; hub direct-pushed to Pi's IP",
         "Register with cmd_port=0 → forces queue delivery; added guard `int(cmd_port) > 0`"),
        ("lora_bridge poller not starting",
         "Poller thread only launched on second registration event",
         "Start poller immediately after first device seen in run_bridge()"),
        ("USB errno -71 (I/O error)",
         "xhci controller in bad state after LoRa32 disconnect",
         "sudo usbreset /dev/bus/usb/001/001 — resets root hub; device re-enumerates"),
        ("Recovery actions all in cooldown",
         "Previous run set 300 s cooldowns; hub not restarted",
         "Reduced cooldown to 60 s for demo; always restart hub between sessions"),
        ("Anomaly detection too sensitive",
         "Baseline contaminated by stress-test data → inflated σ",
         "Use lower 60% of data for baseline; cap at config threshold; baseline freeze"),
        ("Groq LLM rate-limited mid-demo",
         "Free tier: 30 req/min shared across all anomalies",
         "Fallback to Ollama llama3.2:3b (local); Groq only on first anomaly per type"),
        ("Dashboard showing no devices",
         "Hub started with old code before register_local_device() was added",
         "Restart hub with current code; local Pi registers on first metrics push"),
        ("Multiple duplicate processes",
         "SSH restarts created new instances; all racing for /dev/ttyUSB0",
         "pkill -f lora_bridge && fuser -k /dev/ttyUSB0 before each restart"),
    ]

    rh = Emu(330000)
    for i, row in enumerate(rows):
        bg = C_CARD if i % 2 == 0 else C_WHITE
        _table_row(sl, row, col_ws, CT + Emu(320040) + rh * i, rh, bg)


def slide_11_future_work(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Future Work")
    _add_footer(sl, "Current system is a working proof-of-concept; production path is clear")

    gap = Emu(91440)
    col_w = (CW - gap) // 2
    card_h = Emu(620000)

    left_items = [
        ("LoRa Mesh Transport",
         "Replace USB serial with SX1262 radio. Multi-hop mesh across\n"
         "many nodes — no Pi bridge needed for each device."),
        ("Over-the-Air (OTA) Updates",
         "Serve firmware to LoRa32 via LoRa radio or Wi-Fi. "
         "Rolling deploy without physical access."),
        ("Cloud Twin (AWS IoT Core)",
         "Push telemetry to IoT Core → DynamoDB. Shadow documents\n"
         "enable remote config + persistent device state."),
        ("Mobile Alert App",
         "React Native app subscribes to hub WebSocket.\n"
         "Push notifications via FCM/APNs for critical anomalies."),
        ("PCB Design",
         "Replace breadboard with custom 2-layer PCB.\n"
         "Integrate AHT20, OLED, INA219, voltage divider on one board."),
    ]

    right_items = [
        ("LSTM on-device Inference",
         "Quantise LSTM to INT8 TFLite Micro. Run time-series anomaly\n"
         "detection on ESP32-S3 without hub — true edge AI."),
        ("Suricata / Zeek Integration",
         "Connect security agent to real IDS. Replace synthetic threats\n"
         "with actual network traffic analysis."),
        ("Kubernetes Deployment",
         "Containerise hub agents. Horizontal scaling for 100+ devices.\n"
         "Helm chart for one-command deploy."),
        ("Federated Learning",
         "Devices train local Isolation Forest, share only model weights.\n"
         "Privacy-preserving collaborative anomaly learning."),
        ("Energy Harvesting",
         "Solar + supercapacitor for LoRa32. Wake on LoRa interrupt.\n"
         "Target <10 mW average — years on a single charge."),
    ]

    for i, (label, body) in enumerate(left_items):
        top = CT + (card_h + gap) * i
        _add_card(sl, ML, top, col_w, card_h, label, body, body_size=8.5)

    for i, (label, body) in enumerate(right_items):
        top = CT + (card_h + gap) * i
        _add_card(sl, ML + col_w + gap, top, col_w, card_h, label, body, body_size=8.5)


def slide_12_closing(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    _add_header(sl, "Summary")
    _add_footer(sl, "Thank you")

    # Big headline
    _add_textbox(sl, ML, CT + Emu(200000), CW, Emu(600000),
                 "Sentinel AI — end-to-end self-healing IoT infrastructure",
                 font_size_pt=22, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    _add_textbox(sl, ML, CT + Emu(870000), CW, Emu(350000),
                 "From bare-metal sensor reads on ESP32-S3 to AI-driven recovery on a Mac hub",
                 font_size_pt=14, color=C_DARK, align=PP_ALIGN.CENTER)

    # Summary cards
    gap = Emu(91440)
    card_w = (CW - gap * 3) // 4
    card_h = Emu(1600000)
    top = CT + Emu(1300000)

    _add_card(sl, ML, top, card_w, card_h,
              "Sense",
              "AHT20 I²C sensor\nSSD1306 OLED\nCPU / Mem / Disk\nPower (INA219)\nNetwork ping",
              label_size=12, body_size=9)
    _add_card(sl, ML + (card_w + gap), top, card_w, card_h,
              "Detect",
              "z-score · IQR\nTrend elevation\nIsolation Forest\nLSTM autoencoder\nHard floor",
              label_size=12, body_size=9)
    _add_card(sl, ML + (card_w + gap) * 2, top, card_w, card_h,
              "Diagnose",
              "Rule-based YAML\nGroq LLaMA-3.3-70b\nOllama llama3.2:3b\nBackground thread\nFallback chain",
              label_size=12, body_size=9)
    _add_card(sl, ML + (card_w + gap) * 3, top, card_w, card_h,
              "Recover",
              "15+ actions\nGraduated L1–L4\nOutcome check\nPer-device queue\n<30 s pipeline",
              label_size=12, body_size=9)

    # Final quote
    _add_textbox(sl, ML, CT + Emu(3100000), CW, Emu(400000),
                 '"Don\'t wait for things to break — build systems that fix themselves."',
                 font_size_pt=13, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    _add_textbox(sl, ML, CT + Emu(3550000), CW, Emu(300000),
                 "Karthick S  ·  Embedded Systems & IoT  ·  April 2026",
                 font_size_pt=10, color=C_GRAY, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation('/Users/karthi/Downloads/test.pptx')

    # Remove ALL existing slides (work backwards to preserve indices)
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    for sld_id in slide_ids:
        xml_slides.remove(sld_id)

    # Rebuild slides in order
    slide_01_title(prs)
    slide_purpose(prs)
    slide_hardware_photos(prs)
    slide_02_iot_overview(prs)
    slide_03_hardware_summary(prs)
    slide_04_timeline(prs)
    slide_05_sensor_interfacing(prs)
    slide_06_display_output(prs)
    slide_07_stack(prs)
    slide_08_hw_constraints(prs)
    slide_09_sw_constraints(prs)
    slide_10_challenges(prs)
    slide_11_future_work(prs)
    slide_12_closing(prs)

    out = '/Users/karthi/Desktop/Sentinal_AI/Sentinel_AI_Presentation.pptx'
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == '__main__':
    main()

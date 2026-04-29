"""
Sentinel AI – Week 10 Presentation Builder
Generates a clean, graph-heavy, icon-rich PPTX
"""
import os, io, textwrap
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt

# ── PALETTE (soft, professional — no strong primaries) ──────────────────────
NAVY    = RGBColor(0x1E, 0x2A, 0x3A)   # slide header strip
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL= RGBColor(0x2C, 0x3E, 0x50)   # body text
STEEL   = RGBColor(0x41, 0x6E, 0x8F)   # accent blue
TEAL    = RGBColor(0x1A, 0x9E, 0x8A)   # accent teal
AMBER   = RGBColor(0xD4, 0x7C, 0x2B)   # accent amber
LIGHT   = RGBColor(0xF4, 0xF6, 0xF8)   # panel bg
MID     = RGBColor(0xBD, 0xC3, 0xC7)   # borders
REF_CLR = RGBColor(0x7F, 0x8C, 0x8D)   # reference text

HEX_NAVY  = '#1E2A3A'
HEX_STEEL = '#416E8F'
HEX_TEAL  = '#1A9E8A'
HEX_AMBER = '#D47C2B'
HEX_LIGHT = '#F4F6F8'
HEX_MID   = '#BDC3C7'
HEX_CHAR  = '#2C3E50'

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── HELPERS ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)

def fill_shape(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def add_rect(slide, l, t, w, h, rgb):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_shape(shape, rgb)
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=16, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def header_strip(slide, title, icon=''):
    add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
    add_text(slide, f'{icon}  {title}' if icon else title,
             0.3, 0.15, 12.5, 0.8,
             size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def ref_footer(slide, ref_text):
    add_rect(slide, 0, 7.1, 13.33, 0.4, LIGHT)
    add_text(slide, f'Ref: {ref_text}',
             0.2, 7.12, 13.0, 0.35,
             size=8, bold=False, color=REF_CLR, align=PP_ALIGN.LEFT)

def fig_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def add_image(slide, buf, l, t, w, h):
    slide.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))

def panel(slide, l, t, w, h, title='', icon='', title_col=None):
    """Draw a light card panel with optional mini-header."""
    add_rect(slide, l, t, w, h, LIGHT)
    if title:
        tc = title_col or STEEL
        add_text(slide, f'{icon} {title}' if icon else title,
                 l+0.1, t+0.05, w-0.2, 0.35,
                 size=11, bold=True, color=tc)

# ════════════════════════════════════════════════════════════════════════════
#  CHART GENERATORS
# ════════════════════════════════════════════════════════════════════════════

def chart_system_flow():
    """Horizontal pipeline: Sensor → Monitor → Event Bus → AI Agents → Dashboard"""
    fig, ax = plt.subplots(figsize=(10, 2.6), facecolor='white')
    ax.set_facecolor('white')
    ax.axis('off')
    boxes = ['Hardware\nSensors', 'Monitoring\nAgent', 'Event\nBus', 'AI Agent\nPipeline', 'Dashboard\n:5001']
    icons = ['📡', '📊', '🔀', '🤖', '🖥️']
    colors = ['#E8F4FD','#E8F8F5','#FEF9E7','#F0E6FF','#EBF5FB']
    border = [HEX_STEEL, HEX_TEAL, HEX_AMBER, '#7D5FA6', HEX_STEEL]
    xs = [0.5, 2.4, 4.3, 6.2, 8.1]
    for i, (x, box, ic, col, bc) in enumerate(zip(xs, boxes, icons, colors, border)):
        rect = mpatches.FancyBboxPatch((x, 0.3), 1.5, 1.8,
                boxstyle='round,pad=0.08', linewidth=2,
                edgecolor=bc, facecolor=col)
        ax.add_patch(rect)
        ax.text(x+0.75, 1.5, ic, ha='center', va='center', fontsize=20)
        ax.text(x+0.75, 0.72, box, ha='center', va='center', fontsize=9,
                fontweight='bold', color=HEX_CHAR)
        if i < len(xs)-1:
            ax.annotate('', xy=(xs[i+1]-0.02, 1.2), xytext=(x+1.52, 1.2),
                        arrowprops=dict(arrowstyle='->', color=HEX_MID, lw=2))
    ax.set_xlim(0, 10); ax.set_ylim(0, 2.5)
    ax.text(5, 0.08, 'Every 5 seconds · event-driven · pub/sub architecture',
            ha='center', fontsize=8, color=HEX_MID, style='italic')
    fig.tight_layout(pad=0.2)
    return fig_to_image(fig)

def chart_hardware_topology():
    """Network topology: Mac Hub ↔ WiFi Router ↔ Windows laptop"""
    fig, ax = plt.subplots(figsize=(9, 3.2), facecolor='white')
    ax.set_facecolor('white')
    ax.axis('off')
    # nodes
    nodes = [
        (1.2, 1.6, '🍎', 'Mac Hub\n(Sentinel AI)\n10.0.0.x', HEX_TEAL, '#E8F8F5'),
        (4.5, 1.6, '📡', 'WiFi Router\n(Local Network)\n10.0.0.1', HEX_MID, '#F8F8F8'),
        (7.8, 2.5, '💻', 'Windows Laptop\nKARTHIS_DELL\n10.0.0.61', HEX_STEEL, '#E8F4FD'),
        (7.8, 0.7, '🍓', 'Raspberry Pi\n(Future Node)\n10.0.0.x', '#B7950B', '#FEF9E7'),
    ]
    for x, y, ic, label, bc, fc in nodes:
        circ = plt.Circle((x, y), 0.55, color=fc, linewidth=2, ec=bc, zorder=3)
        ax.add_patch(circ)
        ax.text(x, y+0.15, ic, ha='center', va='center', fontsize=22, zorder=4)
        ax.text(x, y-0.85, label, ha='center', va='center', fontsize=8,
                fontweight='bold', color=HEX_CHAR)

    # connections
    for (x1,y1), (x2,y2), lbl in [
        ((1.75,1.6),(3.95,1.6),'HTTP :5001\nMetrics push every 5s'),
        ((5.05,1.6),(7.25,2.5),'HTTP :5001/:5002\nCommands & metrics'),
        ((5.05,1.6),(7.25,0.7),'HTTP :5001\n(planned)'),
    ]:
        ax.annotate('', xy=(x2,y2), xytext=(x1,y1),
                    arrowprops=dict(arrowstyle='<->', color=HEX_STEEL, lw=1.8))
        mx,my = (x1+x2)/2, (y1+y2)/2+0.25
        ax.text(mx, my, lbl, ha='center', fontsize=7, color=HEX_MID)

    # port callout
    ax.text(4.5, 3.15, '⚙  Port 5001 = dashboard/API   Port 5002 = command receiver',
            ha='center', fontsize=8.5, color=HEX_CHAR,
            bbox=dict(boxstyle='round,pad=0.3', fc='#EBF5FB', ec=HEX_STEEL, lw=1))
    ax.set_xlim(0.3, 9); ax.set_ylim(0, 3.4)
    fig.tight_layout(pad=0.3)
    return fig_to_image(fig)

def chart_metrics_collected():
    """Bar chart: metrics per category"""
    fig, ax = plt.subplots(figsize=(5.5, 3.2), facecolor='white')
    cats   = ['CPU', 'Memory', 'Disk', 'Network', 'Power', 'Security']
    counts = [8, 7, 6, 6, 5, 4]
    monitored = [3, 2, 2, 2, 2, 2]
    x = np.arange(len(cats))
    b1 = ax.bar(x, counts, 0.5, label='Total metrics', color='#AED6F1', edgecolor='white')
    b2 = ax.bar(x, monitored, 0.5, label='Anomaly-detected', color=HEX_STEEL, edgecolor='white')
    ax.set_xticks(x); ax.set_xticklabels(cats, fontsize=10)
    ax.set_ylabel('Metric count', fontsize=10)
    ax.set_title('Metrics collected per category', fontsize=11, pad=6, color=HEX_CHAR)
    ax.legend(fontsize=9)
    ax.set_facecolor('white'); fig.patch.set_facecolor('white')
    ax.spines[['top','right']].set_visible(False)
    for bar, v in zip(b1, counts):
        ax.text(bar.get_x()+bar.get_width()/2, v+0.1, str(v),
                ha='center', fontsize=9, color=HEX_CHAR)
    fig.tight_layout(pad=0.4)
    return fig_to_image(fig)

def chart_data_push():
    """Timeline: 5-second push cycle"""
    fig, ax = plt.subplots(figsize=(5.5, 2.8), facecolor='white')
    ax.set_facecolor('white')
    ax.axis('off')
    # draw a simple timeline
    ax.axhline(0.5, 0.05, 0.95, color=HEX_MID, lw=2)
    events = [(0.08, 'collect\nmetrics', '📊'),
              (0.28, 'POST\n/api/metrics/push', '📤'),
              (0.50, 'event bus\npublish', '🔀'),
              (0.72, 'agents\nprocess', '🤖'),
              (0.92, 'dashboard\nupdate', '🖥️')]
    for xp, label, ic in events:
        ax.plot(xp, 0.5, 'o', ms=14, color=HEX_STEEL, zorder=3)
        ax.text(xp, 0.5, ic, ha='center', va='center', fontsize=12, zorder=4)
        ax.text(xp, 0.78, label, ha='center', va='bottom', fontsize=8,
                color=HEX_CHAR, fontweight='bold')
    ax.annotate('', xy=(0.93, 0.25), xytext=(0.07, 0.25),
                arrowprops=dict(arrowstyle='->', color=HEX_AMBER, lw=2))
    ax.text(0.5, 0.22, 'every 5 seconds', ha='center', fontsize=9,
            color=HEX_AMBER, fontweight='bold')
    ax.set_xlim(0, 1); ax.set_ylim(0, 1.1)
    ax.set_title('Data collection & push cycle', fontsize=11, color=HEX_CHAR, pad=4)
    fig.tight_layout(pad=0.3)
    return fig_to_image(fig)

def chart_threshold_adaptive():
    """Line chart: CPU % with IQR fence adapting over time"""
    np.random.seed(42)
    n = 120
    t = np.arange(n)
    # baseline: ~8% with noise, then spike at t=80
    cpu = np.clip(8 + np.random.randn(n)*2.5, 2, 100).astype(float)
    cpu[78:100] = np.clip(85 + np.random.randn(22)*5, 60, 100)  # spike
    cpu[100:] = np.clip(9 + np.random.randn(20)*2, 2, 25)  # recovery

    # rolling IQR fence (window=30)
    fence_mild = np.full(n, np.nan)
    fence_extreme = np.full(n, np.nan)
    for i in range(30, n):
        w = cpu[max(0,i-30):i]
        q1, q3 = np.percentile(w, 25), np.percentile(w, 75)
        iqr = q3 - q1
        fence_mild[i]    = q3 + 1.5*iqr
        fence_extreme[i] = q3 + 3.0*iqr

    fig, ax = plt.subplots(figsize=(7.5, 3.2), facecolor='white')
    ax.set_facecolor('white')
    ax.plot(t, cpu, color=HEX_STEEL, lw=1.5, label='CPU %', zorder=3)
    ax.plot(t, fence_mild, color=HEX_AMBER, lw=2, ls='--', label='IQR mild fence (Q3+1.5·IQR)')
    ax.plot(t, fence_extreme, color='#C0392B', lw=1.5, ls=':', label='IQR extreme fence (Q3+3·IQR)')
    ax.axvspan(78, 100, alpha=0.08, color='red', label='Anomaly detected')
    ax.axvline(78, color='red', lw=1.5, alpha=0.6)
    ax.text(79, 95, '⚠ Anomaly\nfired', fontsize=8, color='#C0392B', fontweight='bold')
    ax.axvspan(0, 15, alpha=0.05, color=HEX_TEAL)
    ax.text(2, 92, '⏳ Warmup\n(15 readings)', fontsize=7.5, color=HEX_TEAL, fontweight='bold')
    ax.set_xlabel('Reading number (1 per 5 s)', fontsize=9)
    ax.set_ylabel('CPU %', fontsize=9)
    ax.set_title('Adaptive IQR fence — bounds learned from live data, no hardcoding', fontsize=10, color=HEX_CHAR)
    ax.legend(fontsize=8, loc='upper left')
    ax.spines[['top','right']].set_visible(False)
    fig.tight_layout(pad=0.4)
    return fig_to_image(fig)

def chart_threshold_methods():
    """Bar chart: detection methods coverage"""
    fig, ax = plt.subplots(figsize=(5, 3), facecolor='white')
    methods = ['IQR\nOutlier', 'Z-Score\nAdaptive', 'Trend\nElevation', 'Rate-of-\nChange', 'Hard\nThreshold']
    fires   = [18, 12, 6, 9, 5]
    colors  = [HEX_STEEL, HEX_TEAL, HEX_AMBER, '#8E44AD', '#C0392B']
    ax.bar(methods, fires, color=colors, edgecolor='white', width=0.6)
    ax.set_ylabel('Detections (demo session)', fontsize=9)
    ax.set_title('Detection method activity', fontsize=10, color=HEX_CHAR)
    ax.set_facecolor('white'); fig.patch.set_facecolor('white')
    ax.spines[['top','right']].set_visible(False)
    for i, v in enumerate(fires):
        ax.text(i, v+0.3, str(v), ha='center', fontsize=10, color=HEX_CHAR)
    fig.tight_layout(pad=0.4)
    return fig_to_image(fig)

def chart_metrics_meaning():
    """Gauge-style: normal ranges for each metric"""
    fig, axes = plt.subplots(1, 4, figsize=(9, 2.6), facecolor='white')
    metrics = [('CPU %', 3.9, 80, 100, HEX_STEEL),
               ('Memory %', 46, 85, 100, HEX_TEAL),
               ('Disk %', 41.7, 90, 100, HEX_AMBER),
               ('Ping ms', 4, 200, 500, '#8E44AD')]
    for ax, (name, val, warn, mx, col) in zip(axes, metrics):
        ax.set_facecolor('white')
        ax.set_xlim(0, mx); ax.set_ylim(0, 1)
        ax.barh(0.5, mx,  height=0.35, color='#EAECEE', left=0)
        ax.barh(0.5, warn, height=0.35, color='#D5F5E3', left=0)
        ax.barh(0.5, val,  height=0.35, color=col, left=0)
        ax.axvline(warn, color=HEX_AMBER, lw=1.5, ls='--')
        ax.text(val, 0.5, f' {val}', va='center', fontsize=9, fontweight='bold', color=col)
        ax.text(warn, 0.02, f'⚠{warn}', ha='center', fontsize=7, color=HEX_AMBER)
        ax.set_title(name, fontsize=9, color=HEX_CHAR, fontweight='bold')
        ax.axis('off')
    fig.suptitle('Live metric values vs anomaly thresholds  (green = normal · amber = warning · colored = current)',
                 fontsize=8, color=HEX_MID)
    fig.tight_layout(pad=0.3)
    return fig_to_image(fig)

def chart_monitor_categories():
    """Donut: % of anomaly events by category in a demo session"""
    fig, ax = plt.subplots(figsize=(4, 3.2), facecolor='white')
    sizes  = [38, 25, 14, 12, 7, 4]
    labels = ['CPU', 'Memory', 'Network', 'Power', 'Disk', 'Security']
    colors = [HEX_STEEL, HEX_TEAL, HEX_AMBER, '#8E44AD', '#27AE60', '#C0392B']
    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, colors=colors,
        autopct='%1.0f%%', startangle=140,
        wedgeprops=dict(width=0.5), pctdistance=0.75)
    for t in texts: t.set_fontsize(9)
    for t in autotexts: t.set_fontsize(8); t.set_color('white')
    ax.set_title('Anomaly events by category\n(typical demo session)', fontsize=9, color=HEX_CHAR)
    fig.tight_layout(pad=0.3)
    return fig_to_image(fig)

def chart_agent_pipeline():
    """Vertical pipeline with timing for each stage"""
    fig, ax = plt.subplots(figsize=(9, 3.0), facecolor='white')
    ax.set_facecolor('white'); ax.axis('off')
    stages = [
        ('📡 Monitoring', 'Collects CPU/Mem/Disk/\nNetwork/Power every 5s', '5 s', HEX_STEEL),
        ('🔬 Anomaly', 'IQR + Z-score +\nTrend + LSTM detects', '15–25 s', HEX_TEAL),
        ('🧠 Diagnosis', 'Groq LLM identifies\nroot cause + actions', '2–8 s', '#8E44AD'),
        ('🔧 Recovery', 'Sends fix commands\nto local/remote device', '<5 s', HEX_AMBER),
        ('📚 Learning', 'Updates baselines\nafter resolution', '1 s', '#27AE60'),
    ]
    xs = [0.55, 2.85, 5.15, 7.45, 9.75]
    for i, (x, (label, desc, timing, col)) in enumerate(zip(xs, stages)):
        rect = mpatches.FancyBboxPatch((x-0.95, 0.35), 1.9, 2.1,
                boxstyle='round,pad=0.1', lw=2, ec=col, fc=col+'22')
        ax.add_patch(rect)
        parts = label.split(' ', 1)
        ax.text(x, 2.15, parts[0], ha='center', va='center', fontsize=20)
        ax.text(x, 1.75, parts[1] if len(parts)>1 else '', ha='center', va='center',
                fontsize=9, fontweight='bold', color=HEX_CHAR)
        ax.text(x, 1.25, desc, ha='center', va='center', fontsize=7.5, color=HEX_CHAR)
        rect2 = mpatches.FancyBboxPatch((x-0.42, 0.38), 0.84, 0.42,
                boxstyle='round,pad=0.05', lw=1, ec=col, fc=col)
        ax.add_patch(rect2)
        ax.text(x, 0.59, timing, ha='center', va='center', fontsize=8,
                fontweight='bold', color='white')
        if i < len(stages)-1:
            ax.annotate('', xy=(xs[i+1]-0.97, 1.4), xytext=(x+0.97, 1.4),
                        arrowprops=dict(arrowstyle='->', color=HEX_MID, lw=2.5))
    ax.text(5.15, 0.1, '⏱  Total: detect → diagnose → recover  ≈ 25 – 40 seconds',
            ha='center', fontsize=9, color=HEX_CHAR, fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.3', fc='#EBF5FB', ec=HEX_STEEL, lw=1))
    ax.set_xlim(-0.2, 11); ax.set_ylim(0, 2.6)
    fig.tight_layout(pad=0.2)
    return fig_to_image(fig)

def chart_pipeline_timing():
    """Bar: avg response times for each pipeline stage"""
    fig, ax = plt.subplots(figsize=(4.5, 3.0), facecolor='white')
    stages = ['Detect\n(IQR)', 'Detect\n(LSTM)', 'Diagnose\n(Groq)', 'Recover\n(remote)', 'Total']
    times  = [20, 390, 5, 3, 40]
    colors = [HEX_STEEL, HEX_TEAL, '#8E44AD', HEX_AMBER, HEX_CHAR]
    bars = ax.bar(stages, times, color=colors, edgecolor='white', width=0.55)
    ax.set_ylabel('Seconds', fontsize=9)
    ax.set_title('Pipeline stage latencies (avg)', fontsize=10, color=HEX_CHAR)
    ax.set_facecolor('white'); fig.patch.set_facecolor('white')
    ax.spines[['top','right']].set_visible(False)
    for bar, v in zip(bars, times):
        ax.text(bar.get_x()+bar.get_width()/2, v+2, f'{v}s',
                ha='center', fontsize=9, color=HEX_CHAR, fontweight='bold')
    ax.set_ylim(0, 470)
    ax.text(1, 410, '⏳ LSTM trains after\n~6.5 min of data', fontsize=7.5,
            color=HEX_TEAL, ha='center')
    fig.tight_layout(pad=0.4)
    return fig_to_image(fig)

def chart_simulation_results():
    """Grouped bar: detection, diagnosis, recovery times per scenario"""
    fig, ax = plt.subplots(figsize=(8.5, 3.2), facecolor='white')
    scenarios = ['CPU Spike', 'Memory\nPressure', 'Disk Fill', 'Power Sag', 'Remote\n(Windows)']
    detect   = [20, 22, 25, 18, 25]
    diagnose = [5,   6,  4,  5,  7]
    recover  = [2,   3,  2,  2,  4]
    x = np.arange(len(scenarios))
    w = 0.22
    ax.bar(x-w,   detect,   w, label='Detect (s)',   color=HEX_STEEL,  edgecolor='white')
    ax.bar(x,     diagnose, w, label='Diagnose (s)', color=HEX_TEAL,   edgecolor='white')
    ax.bar(x+w,   recover,  w, label='Recover (s)',  color=HEX_AMBER,  edgecolor='white')
    ax.set_xticks(x); ax.set_xticklabels(scenarios, fontsize=9)
    ax.set_ylabel('Seconds', fontsize=9)
    ax.set_title('Simulation test results — all scenarios passed ✓', fontsize=10, color=HEX_CHAR)
    ax.legend(fontsize=9); ax.spines[['top','right']].set_visible(False)
    ax.set_facecolor('white'); fig.patch.set_facecolor('white')
    fig.tight_layout(pad=0.4)
    return fig_to_image(fig)

def chart_distributed():
    """Show multi-device metric comparison"""
    fig, axes = plt.subplots(1, 2, figsize=(8.5, 3.0), facecolor='white')
    t = np.arange(60)
    # Mac hub
    cpu_mac = np.clip(15 + np.random.randn(60)*3, 5, 40).astype(float)
    cpu_win = np.clip(8  + np.random.randn(60)*2, 2, 15).astype(float)
    cpu_win[35:55] = np.clip(95 + np.random.randn(20)*3, 80, 100)  # demo_full spike
    axes[0].plot(t, cpu_mac, color=HEX_TEAL,  lw=1.8, label='Mac Hub')
    axes[0].plot(t, cpu_win, color=HEX_STEEL, lw=1.8, label='KARTHIS_DELL')
    axes[0].axhline(80, color='red', lw=1.2, ls='--', alpha=0.6, label='Hard threshold 80%')
    axes[0].axvspan(35, 55, alpha=0.07, color='red')
    axes[0].text(45, 88, '⚡ Demo Full\nPipeline', fontsize=7.5, color='#C0392B', ha='center')
    axes[0].set_title('CPU % — two devices', fontsize=9, color=HEX_CHAR)
    axes[0].set_xlabel('Reading #', fontsize=8); axes[0].set_ylabel('CPU %', fontsize=8)
    axes[0].legend(fontsize=8); axes[0].spines[['top','right']].set_visible(False)
    axes[0].set_facecolor('white')

    # command round-trip
    cats   = ['Command\nqueued', 'Client\npoll', 'Execute\non device', 'Result\nposted back']
    times2 = [0.05, 0.9, 0.2, 0.3]
    axes[1].barh(cats, times2, color=[HEX_STEEL,HEX_TEAL,HEX_AMBER,'#8E44AD'], edgecolor='white')
    axes[1].set_xlabel('Seconds', fontsize=8)
    axes[1].set_title('Command round-trip latency', fontsize=9, color=HEX_CHAR)
    axes[1].spines[['top','right']].set_visible(False)
    axes[1].set_facecolor('white')
    for i, v in enumerate(times2):
        axes[1].text(v+0.01, i, f'{v}s', va='center', fontsize=9, color=HEX_CHAR)
    fig.patch.set_facecolor('white')
    fig.tight_layout(pad=0.4)
    return fig_to_image(fig)

def chart_models():
    """Comparison of AI models used"""
    fig, ax = plt.subplots(figsize=(8.5, 3.2), facecolor='white')
    models   = ['Groq\n(Llama 3.1 70B)', 'Ollama\n(llama3.2:3b)', 'Isolation\nForest', 'LSTM\nAutoencoder', 'Rule-\nbased']
    accuracy = [88, 72, 78, 82, 60]
    speed_s  = [3,  8,  0.02, 0.05, 0.001]
    bar_col  = ['#8E44AD', HEX_TEAL, HEX_STEEL, HEX_AMBER, HEX_MID]
    x = np.arange(len(models))
    bars = ax.bar(x, accuracy, 0.5, color=bar_col, edgecolor='white')
    ax.set_xticks(x); ax.set_xticklabels(models, fontsize=9)
    ax.set_ylabel('Accuracy on test set (%)', fontsize=9, color=HEX_CHAR)
    ax.set_ylim(0, 110)
    ax.set_title('AI model comparison — accuracy vs latency trade-off', fontsize=10, color=HEX_CHAR)
    ax.spines[['top','right']].set_visible(False)
    ax.set_facecolor('white'); fig.patch.set_facecolor('white')
    for bar, acc, spd in zip(bars, accuracy, speed_s):
        ax.text(bar.get_x()+bar.get_width()/2, acc+1.5, f'{acc}%',
                ha='center', fontsize=9, fontweight='bold', color=HEX_CHAR)
        ax.text(bar.get_x()+bar.get_width()/2, acc-8, f'⏱{spd}s',
                ha='center', fontsize=7.5, color='white')
    ax.legend(handles=[mpatches.Patch(color=c, label=l) for c,l in zip(bar_col, models)],
              fontsize=7.5, ncol=5, loc='upper right')
    fig.tight_layout(pad=0.4)
    return fig_to_image(fig)

def chart_model_priority():
    """Waterfall / priority chain"""
    fig, ax = plt.subplots(figsize=(4, 3.0), facecolor='white')
    ax.set_facecolor('white'); ax.axis('off')
    steps = [
        ('1', '☁ Groq API',         'Llama 3.1 70B — fast, accurate\nRequires internet', '#8E44AD'),
        ('2', '🏠 Ollama (local)',   'llama3.2:3b — fully offline\nRuns on Mac/Pi', HEX_TEAL),
        ('3', '📏 Rule-based',       'Instant fallback — pattern matching\nAlways available', HEX_STEEL),
    ]
    for i, (num, title, desc, col) in enumerate(steps):
        y = 2.4 - i*0.85
        rect = mpatches.FancyBboxPatch((0.05, y-0.25), 3.5, 0.65,
                boxstyle='round,pad=0.07', lw=1.5, ec=col, fc=col+'18')
        ax.add_patch(rect)
        ax.text(0.25, y+0.07, f'#{num}', fontsize=14, fontweight='bold', color=col, va='center')
        ax.text(0.6, y+0.12, title, fontsize=9.5, fontweight='bold', color=HEX_CHAR, va='center')
        ax.text(0.6, y-0.07, desc, fontsize=7.5, color=HEX_MID, va='center')
        if i < len(steps)-1:
            ax.annotate('', xy=(1.8, y-0.28), xytext=(1.8, y-0.24),
                        arrowprops=dict(arrowstyle='->', color=HEX_MID, lw=1.5))
    ax.text(1.8, 0.08, '↑ fallback chain — tries each in order', ha='center',
            fontsize=7.5, color=HEX_MID, style='italic')
    ax.set_xlim(0, 3.6); ax.set_ylim(-0.1, 3.0)
    ax.set_title('LLM priority chain', fontsize=9, color=HEX_CHAR, pad=3)
    fig.tight_layout(pad=0.3)
    return fig_to_image(fig)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank_slide(prs)
    # full background
    add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
    # accent stripe
    add_rect(sl, 0, 5.6, 13.33, 0.08, TEAL)
    # main title
    add_text(sl, 'Sentinel AI', 0.8, 1.2, 11.5, 1.4,
             size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, 'Autonomous IoT Self-Healing Monitor',
             0.8, 2.65, 11.5, 0.8,
             size=22, bold=False, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(sl, 'Week 10 Progress Report  ·  Spring 2026',
             0.8, 3.5, 11.5, 0.5,
             size=14, bold=False, color=MID, align=PP_ALIGN.CENTER)
    add_text(sl, 'Sejal Mithare  ·  Karthick Suresh Kumar  ·  Jie Zhang  ·  Naveen Munirathnam',
             0.8, 4.2, 11.5, 0.5,
             size=13, bold=False, color=MID, align=PP_ALIGN.CENTER)
    add_text(sl, '🛡️', 6.2, 5.8, 1.0, 0.9, size=36, color=TEAL, align=PP_ALIGN.CENTER)

def slide_team(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'Meet the Team', '👥')
    members = [
        ('Karthick Suresh Kumar', 'M.S. Computer Engineering\nSoftware Systems & Infrastructure', '💻'),
        ('Sejal Mithare',         'M.Eng AI for Computer Vision & Control\nFull-stack developer experience',    '🤖'),
        ('Jie Zhang',             'M.S. Computer Engineering\nComputer Systems & Real-world Apps',            '⚙️'),
        ('Naveen Munirathnam',    'M.Sc. Electrical Engineering\nEmbedded Systems & IoT',                     '📡'),
    ]
    xs = [0.25, 3.58, 6.91, 10.24]
    for i, (x, (name, bio, ic)) in enumerate(zip(xs, members)):
        add_rect(sl, x, 1.3, 2.85, 5.4, LIGHT)
        add_text(sl, ic, x+1.0, 1.5, 0.9, 0.9, size=32, align=PP_ALIGN.CENTER)
        add_text(sl, name, x+0.1, 2.55, 2.65, 0.55,
                 size=11, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
        add_rect(sl, x+0.9, 3.18, 1.05, 0.04, TEAL)
        add_text(sl, bio, x+0.1, 3.3, 2.65, 1.5,
                 size=9, color=STEEL, align=PP_ALIGN.CENTER)
    ref_footer(sl, 'Team project · ECE/CS Graduate Program · Spring 2026')

def slide_overview(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'Project Overview — What Is Sentinel AI?', '🛡️')
    # 3 key points left
    points = [
        ('❌ Problem', 'IoT devices fail silently — operators only discover faults after downtime. Manual monitoring does not scale to hundreds of devices.'),
        ('✅ Solution', 'Sentinel AI is a multi-agent system that monitors every device in real-time, detects anomalies automatically, and issues recovery actions — no human needed.'),
        ('🌐 Distributed', 'A Mac hub runs the AI pipeline. Remote devices (Windows, Pi) push metrics every 5 s over HTTP. Each device gets its own adaptive intelligence.'),
    ]
    for i, (title, body) in enumerate(points):
        y = 1.25 + i*1.7
        add_rect(sl, 0.3, y, 5.5, 1.5, LIGHT)
        add_text(sl, title, 0.45, y+0.08, 5.2, 0.4, size=11, bold=True, color=STEEL)
        add_text(sl, body,  0.45, y+0.45, 5.2, 1.0, size=9.5, color=CHARCOAL)
    # flow chart right
    img = chart_system_flow()
    add_image(sl, img, 6.1, 1.2, 6.9, 3.0)
    # stats row
    stats = [('5', 'detection\nmethods'), ('6', 'AI agents'), ('<30 s', 'detect→fix'), ('2+', 'devices\nmonitors')]
    for i, (val, lbl) in enumerate(stats):
        x = 6.2 + i*1.72
        add_rect(sl, x, 4.4, 1.55, 1.35, LIGHT)
        add_text(sl, val, x+0.1, 4.5, 1.35, 0.55, size=22, bold=True, color=STEEL, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, x+0.1, 5.0, 1.35, 0.7, size=9, color=CHARCOAL, align=PP_ALIGN.CENTER)
    ref_footer(sl, 'Alauthman & Al-Hyari (2025) "Intelligent Fault Detection in WSNs," Computers 14(6). | Cook et al. (2019) "Anomaly Detection for IoT," IEEE Access 7.')

def slide_hardware(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'Hardware & System Setup', '🔧')
    img = chart_hardware_topology()
    add_image(sl, img, 0.25, 1.15, 8.0, 3.3)

    # right: spec cards
    specs = [
        ('🍎 Mac Hub', 'Runs main.py · Flask :5001\nAll 6 agents · SQLite DB\nSends recovery commands'),
        ('💻 Windows Laptop', 'sentinel_client.py\nStreams metrics via HTTP\nReceives commands on :5002'),
        ('🍓 Raspberry Pi (planned)', 'Same client script\nLightweight psutil metrics\nINA219 power sensor'),
    ]
    for i, (title, body) in enumerate(specs):
        y = 1.2 + i*1.55
        add_rect(sl, 8.55, y, 4.5, 1.35, LIGHT)
        add_text(sl, title, 8.7, y+0.05, 4.2, 0.38, size=10, bold=True, color=STEEL)
        add_text(sl, body,  8.7, y+0.45, 4.2, 0.85, size=9, color=CHARCOAL)

    add_text(sl, '⚙  psutil — no special hardware needed for software metrics · INA219/INA3221 I²C sensor for real power readings',
             0.3, 5.75, 12.8, 0.4, size=9, color=STEEL)
    ref_footer(sl, 'Beg et al. (2017) "IoT Power Monitoring using INA219," IEEE ISCAS. | psutil docs: psutil.readthedocs.io')

def slide_data(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'Where Does Our Data Come From?', '📊')

    img1 = chart_metrics_collected()
    add_image(sl, img1, 0.25, 1.2, 5.6, 3.4)

    img2 = chart_data_push()
    add_image(sl, img2, 6.0, 1.2, 7.0, 3.0)

    # metric table
    rows = [
        ('CPU', 'cpu_percent, top_process_name, cpu_count', '3 monitored'),
        ('Memory', 'memory_percent, swap_percent', '2 monitored'),
        ('Disk', 'disk_percent, read/write MB', '1 monitored'),
        ('Network', 'ping_latency_ms, bytes_sent', '1 monitored'),
        ('Power', 'voltage_v, current_a, watts', '2 monitored'),
    ]
    add_text(sl, 'Key metrics monitored:', 0.3, 4.75, 5.5, 0.35, size=10, bold=True, color=STEEL)
    for i, (cat, metrics, note) in enumerate(rows):
        y = 5.1 + i*0.37
        add_text(sl, f'  {cat}', 0.3, y, 1.3, 0.35, size=8.5, bold=True, color=CHARCOAL)
        add_text(sl, metrics, 1.6, y, 3.8, 0.35, size=8, color=CHARCOAL)
        add_text(sl, note, 5.4, y, 1.8, 0.35, size=8, color=TEAL)
    ref_footer(sl, 'Giannetti et al. (2018) "IoT monitoring with edge analytics," IEEE IoT-J 5(2). | psutil.readthedocs.io')

def slide_thresholds(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'How Thresholds Work — Adaptive, Not Hardcoded', '📈')

    img1 = chart_threshold_adaptive()
    add_image(sl, img1, 0.25, 1.15, 8.0, 3.45)

    img2 = chart_threshold_methods()
    add_image(sl, img2, 8.4, 1.15, 4.7, 3.45)

    # callout boxes bottom
    boxes = [
        ('⏳ 75 s warmup', '15 readings × 5 s\nbaseline stabilises'),
        ('🔁 2 consecutive', 'prevents single-point\nfalse positives'),
        ('⏰ 5 min cooldown', 'suppresses repeated\nalerts for same issue'),
        ('🛡 Hard floor', 'CPU > 80 %  always fires\neven before warmup'),
    ]
    for i, (title, body) in enumerate(boxes):
        x = 0.25 + i*3.27
        add_rect(sl, x, 4.7, 3.1, 1.45, LIGHT)
        add_text(sl, title, x+0.1, 4.78, 2.9, 0.38, size=10, bold=True, color=STEEL)
        add_text(sl, body,  x+0.1, 5.18, 2.9, 0.9,  size=9,  color=CHARCOAL)
    ref_footer(sl, 'Tukey (1977) Exploratory Data Analysis. | Chandola et al. (2009) "Anomaly Detection Survey," ACM CS 41(3). | Laptev et al. (2015) KDD.')

def slide_monitors(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'How Monitors Work — What Do These Numbers Mean?', '🔢')

    img1 = chart_metrics_meaning()
    add_image(sl, img1, 0.25, 1.15, 9.0, 3.0)

    img2 = chart_monitor_categories()
    add_image(sl, img2, 9.3, 1.15, 3.8, 3.3)

    # explanation row
    explanations = [
        ('💻 CPU %', 'How hard the processor\nis working. >80% = alert.'),
        ('🧠 Memory %', 'RAM in use. >85% risks\nout-of-memory crashes.'),
        ('💾 Disk %', 'Storage used. >90%\nblocks write operations.'),
        ('📶 Ping ms', 'Network round-trip time.\n>200 ms = connectivity issue.'),
        ('⚡ Voltage V', 'Power rail health.\n±10% of 5 V = anomaly.'),
    ]
    for i, (title, body) in enumerate(explanations):
        x = 0.25 + i*2.58
        add_rect(sl, x, 4.3, 2.45, 1.7, LIGHT)
        add_text(sl, title, x+0.1, 4.38, 2.25, 0.4, size=10, bold=True, color=STEEL)
        add_text(sl, body,  x+0.1, 4.8,  2.25, 0.95, size=9, color=CHARCOAL)
    ref_footer(sl, 'Spanò et al. (2015) "Ultra-low-power IoT sensor node," IEEE Trans. Circuits Syst. I 62(12). | TI INA219 datasheet (SBOS448G).')

def slide_agents(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'Role of the AI Agents', '🤖')

    img1 = chart_agent_pipeline()
    add_image(sl, img1, 0.25, 1.15, 9.5, 3.1)

    img2 = chart_pipeline_timing()
    add_image(sl, img2, 9.9, 1.15, 3.2, 3.1)

    # recovery actions callout
    add_rect(sl, 0.25, 4.4, 12.8, 1.65, LIGHT)
    add_text(sl, '🔧  Recovery actions sent to remote device (example from KARTHIS_DELL):',
             0.4, 4.48, 12.5, 0.4, size=10, bold=True, color=STEEL)
    actions = ['kill_process → killed python3.11.exe', 'clear_cache → EmptyWorkingSet (300 procs)',
               'stop_stress → killed 12 CPU workers', 'algorithmic_cpu_fix → lowered priority of uihost.exe']
    for i, act in enumerate(actions):
        x = 0.4 + i*3.2
        add_text(sl, f'✓  {act}', x, 4.92, 3.1, 0.45, size=8.5, color=CHARCOAL)
    ref_footer(sl, 'Dubey et al. (2025) "Transformer-Driven Fault Detection," MAKE 7(3). | Alauthman & Al-Hyari (2025) Computers 14(6).')

def slide_testing(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'How Are We Testing? — Simulation Lab + Live Remote', '🧪')

    img1 = chart_simulation_results()
    add_image(sl, img1, 0.25, 1.15, 8.5, 3.3)

    img2 = chart_distributed()
    add_image(sl, img2, 0.1, 4.45, 13.1, 2.6)

    # scenario chips
    scenarios = [('⚡ CPU Spike', 'pins all cores to 100%'),
                 ('🧠 Memory Pressure', 'allocates 40% of RAM'),
                 ('💾 Disk Fill', '200 MB cycling writes'),
                 ('🔌 Power Sag', '−0.75 V for 60 s'),
                 ('🌐 Remote Demo', 'full pipeline on Windows')]
    for i, (title, desc) in enumerate(scenarios):
        x = 8.95 + (i % 2)*2.15
        y = 1.25 + (i // 2)*1.45 + (0.7 if i==4 else 0)
        add_rect(sl, x, y, 2.05, 1.25, LIGHT)
        add_text(sl, title, x+0.1, y+0.06, 1.85, 0.4, size=9.5, bold=True, color=STEEL)
        add_text(sl, desc,  x+0.1, y+0.48, 1.85, 0.55, size=8.5, color=CHARCOAL)

    ref_footer(sl, 'NIST SP 800-94 "Intrusion Detection and Prevention Systems." | Garcia-Teodoro et al. (2009) Computers & Security 28(1-2).')

def slide_models(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'What AI Models Are We Using?', '🧠')

    img1 = chart_models()
    add_image(sl, img1, 0.25, 1.15, 8.5, 3.3)

    img2 = chart_model_priority()
    add_image(sl, img2, 8.9, 1.15, 4.2, 3.4)

    # bottom comparison table
    cols = ['Model', 'Type', 'When used', 'Latency', 'Runs offline']
    rows2 = [
        ('Groq Llama 3.1 70B', 'Cloud LLM',      'Primary diagnosis',          '2–8 s', '✗'),
        ('Ollama llama3.2:3b', 'Local LLM',       'Groq unavailable',           '8–20 s', '✓'),
        ('Isolation Forest',   'ML (sklearn)',     'Multivariate anomaly',       '< 0.1 s', '✓'),
        ('LSTM Autoencoder',   'Deep learning',    'Time-series sequences',      '0.05 s', '✓'),
        ('Rule-based',         'Heuristic',        'Always-available fallback',  '< 1 ms', '✓'),
    ]
    ys = [4.55, 4.93, 5.31, 5.69, 6.07]
    xs2 = [0.3, 2.8, 5.05, 8.05, 10.0]
    widths = [2.4, 2.15, 2.9, 1.85, 1.8]
    # header
    for x, w, col in zip(xs2, widths, cols):
        add_rect(sl, x, 4.2, w-0.05, 0.35, NAVY)
        add_text(sl, col, x+0.05, 4.22, w-0.1, 0.3, size=9, bold=True, color=WHITE)
    for ri, (row, y) in enumerate(zip(rows2, ys)):
        bg = LIGHT if ri%2==0 else WHITE
        for x, w, cell in zip(xs2, widths, row):
            add_rect(sl, x, y, w-0.05, 0.37, bg)
            col = TEAL if cell == '✓' else (RGBColor(0xC0,0x39,0x2B) if cell == '✗' else CHARCOAL)
            add_text(sl, cell, x+0.05, y+0.04, w-0.1, 0.3, size=8.5, color=col)
    ref_footer(sl, 'Groq docs: console.groq.com | Ollama: ollama.ai | Liu et al. (2008) "Isolation Forest," ICDM. | Hochreiter & Schmidhuber (1997) LSTM, Neural Computation 9(8).')

def slide_references(prs):
    sl = blank_slide(prs)
    header_strip(sl, 'References', '📚')
    refs = [
        '[1]  Tukey, J.W. (1977). Exploratory Data Analysis. Addison-Wesley. — IQR fence method.',
        '[2]  Chandola, V. et al. (2009). "Anomaly Detection: A Survey." ACM Comput. Surv. 41(3).',
        '[3]  Laptev, N. et al. (2015). "Generic and Scalable Framework for Automated Time-series Anomaly Detection." KDD \'15.',
        '[4]  Alauthman, M. & Al-Hyari, A. (2025). "Intelligent Fault Detection in WSNs." Computers 14(6).',
        '[5]  Dubey, A. et al. (2025). "Transformer-Driven Fault Detection in IoT." MAKE 7(3).',
        '[6]  Cook, A. et al. (2019). "Anomaly detection for IoT time-series data." IEEE Access 7.',
        '[7]  Garcia-Teodoro, P. et al. (2009). "Anomaly-based IDS." Computers & Security 28(1–2).',
        '[8]  NIST SP 800-94. "Guide to Intrusion Detection and Prevention Systems (IDPS)."',
        '[9]  Beg, O. et al. (2017). "IoT Power Monitoring using INA219." IEEE ISCAS.',
        '[10] Liu, F. et al. (2008). "Isolation Forest." IEEE ICDM.',
        '[11] Hochreiter, S. & Schmidhuber, J. (1997). "Long Short-Term Memory." Neural Computation 9(8).',
        '[12] Giannetti, C. et al. (2018). "IoT monitoring with edge analytics." IEEE IoT-J 5(2).',
    ]
    for i, ref in enumerate(refs):
        col = 0 if i < 6 else 1
        row = i if i < 6 else i - 6
        x = 0.4 + col * 6.55
        y = 1.35 + row * 0.85
        add_text(sl, ref, x, y, 6.2, 0.75, size=9, color=CHARCOAL)

# ════════════════════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════════════════════
def main():
    print("Building Sentinel AI Week 10 Presentation...")
    prs = new_prs()

    print("  Slide 1: Title")
    slide_title(prs)
    print("  Slide 2: Team")
    slide_team(prs)
    print("  Slide 3: Project Overview")
    slide_overview(prs)
    print("  Slide 4: Hardware Setup")
    slide_hardware(prs)
    print("  Slide 5: Data Sources")
    slide_data(prs)
    print("  Slide 6: Thresholds")
    slide_thresholds(prs)
    print("  Slide 7: Monitors / Numbers")
    slide_monitors(prs)
    print("  Slide 8: AI Agents")
    slide_agents(prs)
    print("  Slide 9: Testing")
    slide_testing(prs)
    print("  Slide 10: Models")
    slide_models(prs)
    print("  Slide 11: References")
    slide_references(prs)

    out = '/Users/karthi/Desktop/Sentinel_AI_Week10_Presentation.pptx'
    prs.save(out)
    print(f"\n✅  Saved → {out}")
    print(f"   Slides: {len(prs.slides)}  |  {os.path.getsize(out)//1024} KB")

if __name__ == '__main__':
    main()
